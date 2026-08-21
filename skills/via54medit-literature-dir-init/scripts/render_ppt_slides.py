#!/usr/bin/env python3
"""统一 PPT 渲染 SOP — 双引擎 (PowerPoint AppleScript + LibreOffice)

核心原则: 输出到 PPT 同目录树下的 _ppt_renders/ (避免沙盒授权)

用法:
    python3 render_ppt_slides.py <pptx_path> [--engine applescript|libreoffice]

输出:
    <ppt_dir>/_ppt_renders/
    ├── slide_001.jpg ~ slide_NNN.jpg (150 DPI)
    ├── citation_table.csv (A_slide, B_mark, C_citation, D_ppt_content_visual_text)
    └── render_manifest.json
"""
import os, sys, json, subprocess, re, csv
from collections import defaultdict
from pptx import Presentation


def render_ppt_applescript(pptx_path, output_dir):
    """PowerPoint AppleScript 引擎: 输出到 PPT 同目录树 (避免授权)"""
    pdf_path = os.path.join(output_dir, '_ppt_export.pdf')
    if os.path.exists(pdf_path):
        os.remove(pdf_path)
    ascript = f'''
    tell application "Microsoft PowerPoint"
        open POSIX file "{pptx_path}"
        delay 5
        set theDoc to active presentation
        save theDoc in POSIX file "{pdf_path}" as save as PDF
        delay 2
        close theDoc saving no
    end tell
    '''
    r = subprocess.run(['osascript', '-e', ascript], capture_output=True, text=True, timeout=120)
    if r.returncode != 0:
        raise RuntimeError(f'PowerPoint AppleScript 失败: {r.stderr[:500]}')
    if not os.path.exists(pdf_path):
        raise RuntimeError('PowerPoint 导出 PDF 未生成')
    return pdf_path


def render_ppt_libreoffice(pptx_path, output_dir):
    """LibreOffice 引擎 (备选)"""
    pdf_path = os.path.join(output_dir, '_ppt_export.pdf')
    cmd = ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path]
    r = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
    if r.returncode != 0:
        raise RuntimeError(f'LibreOffice 失败: {r.stderr[:500]}')
    expected = os.path.join(output_dir, os.path.basename(pptx_path).replace('.pptx', '.pdf'))
    if os.path.exists(expected) and expected != pdf_path:
        os.rename(expected, pdf_path)
    return pdf_path


def pdf_to_jpgs(pdf_path, output_dir):
    """fitz PDF → JPG (150 DPI)"""
    import fitz
    doc = fitz.open(pdf_path)
    slides = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        jpg = os.path.join(output_dir, f'slide_{i+1:03d}.jpg')
        pix.save(jpg)
        slides.append({'slide_num': i+1, 'jpg': jpg, 'w': pix.width, 'h': pix.height})
    doc.close()
    return slides


def extract_citations(pptx_path):
    """提取全部参考文献 (语义 + 序号智能切分)

    支持两种格式:
    - 多段落: 每条引用换行 (\\n 自然分隔)
    - 单段落: 多条引用挤在一起, 按 "数字. 作者/期刊" 语义位置切分 (finditer)
    """
    prs = Presentation(pptx_path)
    SLIDE_H = prs.slide_height
    all_refs = []
    seen_pairs = set()

    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not hasattr(shape, 'text') or not shape.text.strip():
                continue
            if shape.top < SLIDE_H * 0.70:
                continue
            t = shape.text.strip()
            if not re.search(r'\d{1,3}\.\s*[A-Z\u4e00-\u9fff]', t):
                continue

            for line in [l.strip() for l in t.split('\n') if l.strip()]:
                matches = list(re.finditer(
                    r'(\d{1,3})\.\s*(?=[A-Z\u4e00-\u9fff\u00c0-\u024f])', line
                ))
                for i, m in enumerate(matches):
                    num = int(m.group(1))
                    start = m.end()
                    end = matches[i+1].start() if i+1 < len(matches) else len(line)
                    text = line[start:end].strip()
                    text = re.sub(r'\s+', ' ', text)[:300].rstrip(' ,.,')
                    if len(text) < 10:
                        continue
                    pair = (idx, num)
                    if pair not in seen_pairs:
                        seen_pairs.add(pair)
                        all_refs.append({'slide': idx, 'num': num, 'text': text})

    return all_refs


def write_citation_table(unique_refs, output_dir):
    csv_path = os.path.join(output_dir, 'citation_table.csv')
    with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
        w = csv.writer(f)
        w.writerow(['A_slide', 'B_mark', 'C_citation', 'D_ppt_content_visual_text'])
        for r in unique_refs:
            w.writerow([r['slide'], r['num'], r['text'], ''])
    return csv_path


def render(pptx_path, output_dir, engine='applescript'):
    """统一入口"""
    os.makedirs(output_dir, exist_ok=True)

    print(f'[引擎] {engine}')
    if engine == 'applescript':
        pdf_path = render_ppt_applescript(pptx_path, output_dir)
    else:
        pdf_path = render_ppt_libreoffice(pptx_path, output_dir)
    print(f'[PDF]  {os.path.getsize(pdf_path):,} bytes')

    slides = pdf_to_jpgs(pdf_path, output_dir)
    print(f'[JPG]  {len(slides)} 张 ({slides[0]["w"]}x{slides[0]["h"]})')

    refs = extract_citations(pptx_path)
    csv_path = write_citation_table(refs, output_dir)
    print(f'[CSV]  {len(refs)} 条引用 → {csv_path}')

    manifest = {
        'pptx': pptx_path, 'engine': engine,
        'total_slides': len(slides), 'total_refs': len(refs),
        'output_dir': output_dir,
        'slides': [s['slide_num'] for s in slides],
        'refs': [{'num': r['num'], 'slide': r['slide'], 'text': r['text'][:60]} for r in refs],
    }
    with open(os.path.join(output_dir, 'render_manifest.json'), 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    return manifest


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='PPT 渲染（相对目录 + 双引擎）')
    parser.add_argument('pptx', help='PPT 路径')
    parser.add_argument('--engine', choices=['applescript', 'libreoffice'], default='applescript',
                       help='渲染引擎 (默认 applescript, 避免授权)')
    args = parser.parse_args()

    ppt_dir = os.path.dirname(os.path.abspath(args.pptx))
    output_dir = os.path.join(ppt_dir, '_ppt_renders')

    print(f'🎯 渲染: {os.path.basename(args.pptx)}')
    print(f'   输出: {output_dir}/ (同目录树, 无需授权)')

    manifest = render(args.pptx, output_dir, engine=args.engine)
    print(f'✅ 完成: {manifest["total_slides"]} slides, {manifest["total_refs"]} refs')
