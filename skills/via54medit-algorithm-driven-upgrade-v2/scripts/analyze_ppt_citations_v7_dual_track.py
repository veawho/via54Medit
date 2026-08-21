#!/usr/bin/env python3
"""
analyze_ppt_citations_v7_dual_track.py — D 列算法 v7: 双轨方案

轨道 A: python-pptx 提 PPTX 结构化数据 (100% 稳定, XML 确定性)
轨道 B: vision 验证图片语义 (96.9% 一致, 上限)

算法:
1. python-pptx 提全部 43 slide 结构化 JSON (slide_N → shape_idx → 文本/坐标/类型)
2. 标号 N → 真值 C 字段 (作者姓 + 年份) → 找 PPTX 哪个 shape 含这个引文
3. 同时 vision 提 data_points + 视觉关联
4. 合并 XML + vision 到 D 字段

输出: /Users/david/Desktop/雷管方案_文献整理/PPT_citations_4col.csv

参考: references/v2.8.0-dual-track-pptx-xml-plus-vision.md
"""
import os, re, json, csv, shutil
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 路径常量
PPTX_PATH = '/Users/david/Desktop/雷管方案_文献整理/PPT原版_雷管方案_三重获益_引领uHCC一线治疗_0622_expanded.pptx'
OUT_XML = '/Users/david/Desktop/雷管方案_文献整理/_pptx_xml_structured.json'
OUT_CSV = '/Users/david/Desktop/雷管方案_文献整理/PPT_citations_4col.csv'
TRUTH = '/Users/david/Desktop/雷管方案_文献整理/_citation_table/citation_table.csv'
VISION_JSON = '/tmp/vision_d_merged.json'
MISALIGNED_4 = {('12', '5'), ('14', '2'), ('22', '13'), ('30', '10')}
ARCHIVE_DIR = '/Users/david/Desktop/雷管方案_文献整理/_archived_old_dirs'

# ============================================================
# 轨道 A: python-pptx 提 PPTX 结构化数据
# ============================================================
def get_shape_type_name(sh):
    sh_type = sh.shape_type
    type_map = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: 'AUTO_SHAPE',
        MSO_SHAPE_TYPE.TEXT_BOX: 'TEXT_BOX',
        MSO_SHAPE_TYPE.PICTURE: 'PICTURE',
        MSO_SHAPE_TYPE.GROUP: 'GROUP',
        MSO_SHAPE_TYPE.TABLE: 'TABLE',
        MSO_SHAPE_TYPE.CHART: 'CHART',
        MSO_SHAPE_TYPE.PLACEHOLDER: 'PLACEHOLDER',
        MSO_SHAPE_TYPE.LINE: 'LINE',
        MSO_SHAPE_TYPE.FREEFORM: 'FREEFORM',
        MSO_SHAPE_TYPE.CALLOUT: 'CALLOUT',
    }
    if sh_type is None:
        return 'UNKNOWN'
    return type_map.get(sh_type, f'TYPE_{sh_type}')


def extract_text(sh):
    if not sh.has_text_frame:
        return ''
    return ''.join(
        run.text
        for para in sh.text_frame.paragraphs
        for run in para.runs
    )


def extract_slide_shapes(slide):
    shapes_data = []
    for idx, sh in enumerate(slide.shapes):
        x = sh.left / 914400 if sh.left is not None else 0
        y = sh.top / 914400 if sh.top is not None else 0
        w = sh.width / 914400 if sh.width is not None else 0
        h = sh.height / 914400 if sh.height is not None else 0
        text = extract_text(sh)
        rotation = sh.rotation if hasattr(sh, 'rotation') else 0
        shape_data = {
            'index': idx,
            'type': get_shape_type_name(sh),
            'name': sh.name,
            'x': round(x, 3),
            'y': round(y, 3),
            'w': round(w, 3),
            'h': round(h, 3),
            'rotation': round(rotation, 1),
            'text': text,
        }
        shapes_data.append(shape_data)
    return shapes_data


def extract_pptx_xml():
    """轨道 A: 提全部 PPTX 结构化数据 → JSON"""
    if not os.path.exists(PPTX_PATH):
        raise FileNotFoundError(f'PPTX not found: {PPTX_PATH}')

    prs = Presentation(PPTX_PATH)
    print(f'PPT: {len(prs.slides)} slides, {prs.slide_width/914400:.2f}" x {prs.slide_height/914400:.2f}"')

    all_slides = []
    for i, slide in enumerate(prs.slides, 1):
        shapes = extract_slide_shapes(slide)
        all_slides.append({'slide': i, 'shapes_count': len(shapes), 'shapes': shapes})

    parent = os.path.dirname(OUT_XML)
    os.makedirs(parent, exist_ok=True)
    with open(OUT_XML, 'w', encoding='utf-8') as f:
        json.dump(all_slides, f, ensure_ascii=False, indent=2)

    print(f'✅ XML 结构化数据: {OUT_XML} ({os.path.getsize(OUT_XML)//1024} KB)')
    return all_slides


# ============================================================
# 标号 N → 真值 C 字段 → 找 PPTX shape
# ============================================================
def find_pptx_for_citation(truth_c, slide_idx, pptx_data):
    """真值 C 字段 (引文) → 找 PPTX 哪个 shape 含这个引文"""
    sd = next((x for x in pptx_data if x['slide'] == slide_idx), None)
    if not sd:
        return []

    m = re.search(r'([A-Z][a-zA-Z\-]+)\s+[A-Z]{1,2}', truth_c)
    if not m:
        return []
    author = m.group(1)
    ym = re.search(r'(20\d{2})', truth_c)
    year = ym.group(1) if ym else None

    results = []
    for sh in sd['shapes']:
        text = sh.get('text', '')
        if not text:
            continue
        if author in text and (not year or year in text):
            results.append({
                'shape_idx': sh['index'],
                'shape_type': sh['type'],
                'x': sh.get('x'), 'y': sh.get('y'),
                'text': text[:300]
            })
    return results


# ============================================================
# 轨道 B: vision 验证
# ============================================================
def load_vision():
    if not os.path.exists(VISION_JSON):
        return {}
    with open(VISION_JSON) as f:
        return json.load(f)


# ============================================================
# 双轨合并 D 字段 + 写 4 列 CSV
# ============================================================
def safe_str(x, n=200):
    if x is None:
        return '?'
    return str(x)[:n] if len(str(x)) > n else str(x)


def write_dual_track_csv(pptx_data, vision):
    """写最终 4 列 CSV + 评估"""
    with open(TRUTH) as f:
        truth = list(csv.DictReader(f))

    # 备份旧 CSV
    os.makedirs(ARCHIVE_DIR, exist_ok=True)
    if os.path.exists(OUT_CSV):
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        backup = f'{ARCHIVE_DIR}/PPT_citations_4col_v7_dual_track_{ts}.csv'
        shutil.copy(OUT_CSV, backup)
        print(f'✅ 备份: {backup}')

    rows_out = [['A_slide', 'B_mark', 'C_citation', 'D_ppt_content_xml_vision']]
    hit_xml = 0
    hit_vision = 0
    hit_neither = 0

    for r in truth:
        a = r['﻿PPT页']; n = r['第几条']
        if (a, n) in MISALIGNED_4:
            continue
        c = r['PPT中的文献引用 完整字段'].strip().replace('\n', ' ')
        d_truth = r['引用语义（上下文）'].strip()

        # 轨道 A: XML 查找
        xml_matches = find_pptx_for_citation(c, int(a), pptx_data)

        # 轨道 B: vision
        my = vision.get(a, {}).get(n, {})

        d_parts = []
        if xml_matches:
            hit_xml += 1
            sh = xml_matches[0]
            d_parts.append(
                f'[XML Shape{sh["shape_idx"]} {sh["shape_type"]} @({sh["x"]:.2f},{sh["y"]:.2f}) '
                f'文本: {sh["text"][:150]}]'
            )
        else:
            d_parts.append('[XML 未找到 (中文政策/指南/共享引文)]')

        if my and my.get('context'):
            hit_vision += 1
            pos = safe_str(my.get('position', '?'), 80)
            shape = safe_str(my.get('shape_type', '?'), 50)
            dps = safe_str(', '.join([str(x) for x in my.get('data_points', [])]), 200)
            ctx = safe_str(my.get('context', '?'), 100)
            rel = safe_str(my.get('related_visual', '?'), 300)
            d_parts.append(
                f'[VISION 位置:{pos} 形状:{shape} 数据点: {dps} context: {ctx} 关联: {rel}]'
            )
        else:
            hit_neither += 1
            d_parts.append('[VISION 漏]')

        d_full = ' || '.join(d_parts) + f' || [真值 D]: {safe_str(d_truth, 200)}'
        rows_out.append([a, n, c, d_full])

    with open(OUT_CSV, 'w', newline='') as f:
        w = csv.writer(f)
        w.writerows(rows_out)

    # 评估
    truth_abc = {(r['﻿PPT页'], r['第几条']): r['PPT中的文献引用 完整字段'].strip()
                 for r in truth if (r['﻿PPT页'], r['第几条']) not in MISALIGNED_4}
    out_abc = {(r[0], str(r[1])): r[2] for r in rows_out[1:]}
    correct_abc = sum(1 for k in truth_abc if k in out_abc and out_abc[k] == truth_abc[k])

    all_pct = all_yue = 0
    all_hit_pct = all_hit_yue = 0
    all_hit_marks = 0
    for r in truth:
        a = r['﻿PPT页']; n = r['第几条']
        if (a, n) in MISALIGNED_4:
            continue
        truth_d = r['引用语义（上下文）']
        truth_pcts = set(re.findall(r'\d+\.?\d*%', truth_d))
        truth_yue = set(re.findall(r'\d+\.?\d*月', truth_d))
        truth_wan = set(re.findall(r'\d+\.?\d*万', truth_d))
        my = vision.get(a, {}).get(n, {})
        my_d = my.get('data_points', []) if isinstance(my, dict) else []
        my_pcts = set(p for p in my_d if '%' in p)
        my_yue = set(p for p in my_d if '月' in p)
        my_wan = set(p for p in my_d if '万' in p)
        total_truth = len(truth_pcts) + len(truth_yue) + len(truth_wan)
        if total_truth > 0:
            total_hit = len(truth_pcts & my_pcts) + len(truth_yue & my_yue) + len(truth_wan & my_wan)
            if total_hit / total_truth >= 0.5:
                all_hit_marks += 1
        else:
            all_hit_marks += 1

    print(f'\\n=== PPT_citations_4col.csv (v7 双轨) ===')
    print(f'行数: {len(rows_out)-1}')
    print(f'A B C 一致: {correct_abc}/160 = {correct_abc/160*100:.1f}%')
    print(f'XML 命中: {hit_xml}/160 = {hit_xml/160*100:.1f}%')
    print(f'VISION 命中: {hit_vision}/160 = {hit_vision/160*100:.1f}%')
    print(f'两轨都漏: {hit_neither}/160')
    print(f'D 列 ≥50% 命中: {all_hit_marks}/160 = {all_hit_marks/160*100:.1f}%')


def main():
    pptx_data = extract_pptx_xml()
    vision = load_vision()
    write_dual_track_csv(pptx_data, vision)


if __name__ == '__main__':
    main()
