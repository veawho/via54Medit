#!/usr/bin/env python3
"""
tma_highlight_by_slide.py — 按 slide 顺序驱动的应证 highlight (v3 FINAL rect + 表格 + 图表/图片)

流程 (用户规范):
  0) 提前导出全部 PPT slide 图 → _ppt_renders/slide_NNN.png
     (python-pptx + Pillow 近似渲染; 环境有 LibreOffice/PowerPoint 时可用真实渲染替换)
  1) 逐 slide 视觉提取: 文本块 / 表格 cell / 图片形状 (python-pptx, 融合 _vision_report.json)
  2) 收集该 slide 的所有 PDF (_2_pdfs/P{slide}-*.pdf)
  3) 对每个 PDF 找印证 PPT 内容的全部内容并 highlight:
       a. 文字段落 → 语义匹配 → 逐行 rect (v3 FINAL, RGB 255,217,0, opacity 0.45)
       b. 表格 → PyMuPDF find_tables → 表格 bbox rect (表格文本命中 PPT 术语/数据点)
       c. 图表/图片 → PyMuPDF get_image_info → 图片 bbox rect (该页已有应证命中)
  4) 9 铁律 (仅文字 rect) + 导出高亮页图 + 全部页图
输出: {out-dir}/P{slide}-{num}/  (标准嵌套: main/highlight pdf + 高亮页图 + pages/)

用法:
  python tma_highlight_by_slide.py [--project-dir DIR] [--slide N] [--all] [--force] [--no-render] [--only P3-1]
"""
import os, re, sys, io, json, argparse, shutil
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "hl_v3_final"))

import fitz
from hl_lib import highlight_sentences
from via54_ppt_visual_to_pdf import find_pdf_visual_match

# ============ TMA 核心术语表 (中英) ============
TERMS = [
    "tma", "ttp", "hus", "ahus", "pnh", "hct", "hsct", "gvhd",
    "complement", "补体", "endothelial", "内皮", "microangiopathy", "微血管病",
    "thrombotic", "血栓", "hemolytic", "溶血", "uremic", "尿毒",
    "thrombocytopenia", "血小板减少", "schistocyte", "破碎红细胞",
    "transplant", "移植", "conditioning", "预处理", "calcineurin", "他克莫司",
    "sirolimus", "eculizumab", "依库珠单抗", "ravulizumab", "c5", "c3",
    "mac", "membrane attack", "膜攻击", "proteinuria", "蛋白尿",
    "hypertension", "高血压", "lactate", "乳酸", "ldh", "bilirubin", "胆红素",
    "anemia", "贫血", "platelet", "血小板", "creatinine", "肌酐",
    "drug-induced", "药物相关", "infection", "感染", "autoimmune", "自身免疫",
    "malignancy", "肿瘤", "pregnancy", "妊娠", "narsoplimas", "defibrotide",
    "去纤苷", "rituximab", "利妥昔单抗", "plasma exchange", "血浆置换",
    "therapeutic plasma", "tpe", "immunosuppression", "免疫抑制",
]
DATA_PAT = re.compile(r"\d{1,3}(?:\.\d+)?\s*(?:%|月|天|年)|n\s*=\s*\d+|OR\s*[\d.]+|HR\s*[\d.]+|\b\d{1,2}(?:\.\d+)?\s*mg")


def project_root(ns):
    if getattr(ns, "project_dir", None):
        return ns.project_dir
    return os.environ.get("TMA_PROJECT") or ""


# ============ Step 0: 导出 PPT slide 图 ============
def render_ppt_slides(pptx_path, out_dir):
    """python-pptx + Pillow 近似渲染每页 slide 为 PNG (存档/视觉核对用)"""
    os.makedirs(out_dir, exist_ok=True)
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        print("  [warn] render 依赖缺失: %s (跳过 slide 图导出)" % e)
        return 0
    font_path = None
    for cand in [r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\msyhl.ttc", r"C:\Windows\Fonts\simhei.ttf"]:
        if os.path.exists(cand):
            font_path = cand
            break
    prs = Presentation(pptx_path)
    EMU_IN = 914400.0
    DPI = 120
    n = 0
    for idx, slide in enumerate(prs.slides, start=1):
        w = int(prs.slide_width / EMU_IN * DPI)
        h = int(prs.slide_height / EMU_IN * DPI)
        img = Image.new("RGB", (max(w, 1), max(h, 1)), "white")
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    bio = io.BytesIO(shape.image.blob)
                    im = Image.open(bio).convert("RGB")
                    x0 = int(shape.left / EMU_IN * DPI); y0 = int(shape.top / EMU_IN * DPI)
                    x1 = int((shape.left + shape.width) / EMU_IN * DPI)
                    y1 = int((shape.top + shape.height) / EMU_IN * DPI)
                    if x1 > x0 and y1 > y0:
                        img.paste(im.resize((x1 - x0, y1 - y0)), (x0, y0))
                elif shape.has_table:
                    tbl = shape.table
                    x0 = int(shape.left / EMU_IN * DPI); y0 = int(shape.top / EMU_IN * DPI)
                    x1 = int((shape.left + shape.width) / EMU_IN * DPI)
                    y1 = int((shape.top + shape.height) / EMU_IN * DPI)
                    draw.rectangle([x0, y0, x1, y1], outline="black")
                    rows = len(tbl.rows); cols = len(tbl.columns)
                    rh = (y1 - y0) / rows if rows else 0
                    cw = (x1 - x0) / cols if cols else 0
                    for ri in range(rows):
                        for ci in range(cols):
                            cell = tbl.cell(ri, ci)
                            txt = (cell.text or "").strip()[:40]
                            cx0 = x0 + ci * cw; cy0 = y0 + ri * rh
                            draw.rectangle([cx0, cy0, cx0 + cw, cy0 + rh], outline="black")
                            if txt:
                                fnt = ImageFont.truetype(font_path, max(int(rh * 0.5), 8)) if font_path else ImageFont.load_default()
                                draw.text((cx0 + 2, cy0 + 2), txt, fill="black", font=fnt)
                elif shape.has_text_frame:
                    txt = (shape.text_frame.text or "").strip()
                    if not txt:
                        continue
                    x0 = int(shape.left / EMU_IN * DPI); y0 = int(shape.top / EMU_IN * DPI)
                    fnt = ImageFont.truetype(font_path, 12) if font_path else ImageFont.load_default()
                    draw.text((x0 + 2, y0 + 2), txt[:200], fill="black", font=fnt)
            except Exception:
                continue
        out = os.path.join(out_dir, "slide_%03d.png" % idx)
        img.save(out)
        n += 1
    return n


# ============ Step 1: slide 视觉提取 ============
def extract_slide_visual(pptx_path, slide_num, vision=None):
    """python-pptx 提取 slide 文本块/表格/图片 + 融合 _vision_report.json"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num - 1]
    text_blocks, tables, images = [], [], []
    for shape in slide.shapes:
        try:
            if shape.has_table:
                cells = []
                for r in shape.table.rows:
                    for c in r.cells:
                        if c.text and c.text.strip():
                            cells.append(c.text.strip())
                tables.append({"text": " | ".join(cells)[:500], "n_cells": len(cells)})
            elif getattr(shape, "image", None) is not None:
                images.append({"has_image": True})
            elif shape.has_text_frame:
                t = shape.text_frame.text or ""
                if t.strip():
                    text_blocks.append({"text": t[:300], "position": "center"})
        except Exception:
            continue
    visual = {
        "slide_num": slide_num,
        "text_blocks": text_blocks,
        "citation_marks": [],
        "data_points": [],
        "tables": tables,
        "images": images,
    }
    # 融合 vision report (若含该页)
    if vision:
        sv = vision.get("slides", {}).get(str(slide_num))
        if sv:
            for mid, mark in (sv.get("citation_marks") or {}).items():
                ctx = (mark.get("context") or "").strip()
                if ctx:
                    visual["citation_marks"].append({"mark": mid, "context_text": ctx, "visual_position": "center"})
            for dp in sv.get("data_points") or []:
                if isinstance(dp, dict) and dp.get("context"):
                    visual["data_points"].append(dp)
    # 数据点兜底: 从文本块提取百分比/数字
    for tb in text_blocks:
        for m in DATA_PAT.finditer(tb["text"]):
            visual["data_points"].append({"text": m.group(0), "type": "number", "context": tb["text"][:120]})
    return visual


def slide_terms(visual):
    """slide 关键词集合 = 术语表命中 + 数据点"""
    terms = set(TERMS)
    data = []
    for dp in visual.get("data_points") or []:
        txt = (dp.get("text") or dp.get("context") or "")
        data.append(txt)
    return terms, data


# ============ Step 3: 表格 / 图片应证 ============
def find_table_matches(page, terms, data):
    """返回 [(bbox, hits), ...] 表格命中 (表格文本含 >=2 术语 或 >=1 数据点)"""
    out = []
    try:
        finder = page.find_tables()
    except Exception:
        return out
    for t in finder.tables:
        try:
            cells = t.extract()
            flat = " ".join(str(c or "") for row in cells for c in row).lower()
            if not flat.strip():
                continue
            hits = sum(1 for term in terms if term.lower() in flat)
            dhit = sum(1 for d in data if d and d.lower() in flat)
            if hits >= 2 or (hits >= 1 and dhit >= 1):
                out.append((t.bbox, hits + dhit))
        except Exception:
            continue
    return out


def find_image_matches(page, terms, data, page_has_hit):
    """该页已有文字/表格命中时, 高亮面积 >=4% 页面的图片"""
    out = []
    if not page_has_hit:
        return out
    try:
        infos = page.get_image_info()
    except Exception:
        return out
    page_area = max(page.rect.get_area(), 1e-6)
    for info in infos:
        bbox = info.get("bbox")
        if not bbox:
            continue
        r = fitz.Rect(bbox)
        if r.get_area() / page_area >= 0.04:
            out.append((r, "image"))
    return out[:3]  # 每页最多 3 个图片 rect


# ============ 主流程 ============
def highlight_one_pdf(pdf_path, slide_visual, out_dir, apply_9_rules=True):
    """单个 PDF: 文字(语义) + 表格 + 图片 应证 highlight"""
    pdf_basename = os.path.basename(pdf_path).replace(".pdf", "")
    pn_x = pdf_basename
    m_old = re.match(r"Pn-S(\d+)_(\d+)$", pn_x)
    if m_old:
        pn_x = "P%s-%s" % (m_old.group(1), m_old.group(2))
    os.makedirs(out_dir, exist_ok=True)
    terms, data = slide_terms(slide_visual)

    # --- 文字应证 (语义匹配) ---
    sentences_map = find_pdf_visual_match(pdf_path, slide_visual)
    hl_tmp = os.path.join(out_dir, pn_x + "_highlight_tmp.pdf")
    report = highlight_sentences(pdf_path, hl_tmp, sentences_map, verbose=False)

    # --- 表格 / 图片应证 ---
    doc = fitz.open(hl_tmp)
    keep_rects = []
    for pi in range(len(doc)):
        page = doc[pi]
        page_hit = bool(sentences_map.get(pi))
        for bbox, _hits in find_table_matches(page, terms, data):
            r = fitz.Rect(bbox)
            if r.get_area() <= 0:
                continue
            annot = page.add_rect_annot(r)
            annot.set_colors(stroke=(1.0, 0.85, 0.0), fill=(1.0, 0.85, 0.0))
            annot.set_border(width=0)
            annot.set_opacity(0.45)
            annot.update()
            keep_rects.append((pi, r.round()))
        for r, _kind in find_image_matches(page, terms, data, page_hit):
            annot = page.add_rect_annot(r)
            annot.set_colors(stroke=(1.0, 0.85, 0.0), fill=(1.0, 0.85, 0.0))
            annot.set_border(width=0)
            annot.set_opacity(0.45)
            annot.update()
            keep_rects.append((pi, r.round()))

    # --- 9 铁律 (仅文字 rect; 表格/图片区域跳过) ---
    # 注: PyMuPDF 1.28.2 的 annot.rect 访问存在原生崩溃 bug (0xC0000005),
    #     统一走 xref 层读 Rect, annot API 仅用于枚举/删除。
    removed = 0
    if apply_9_rules:
        from via54_highlight_v3_final import is_metadata_rect
        for pi in range(len(doc)):
            page = doc[pi]
            for annot in list(page.annots() or []):
                xref = annot.xref
                rv = doc.xref_get_key(xref, "Rect")[1]
                if not rv.startswith("["):
                    continue
                try:
                    nums = [float(x) for x in rv.strip("[]").split()]
                    rect = fitz.Rect(nums)
                except Exception:
                    continue
                if not rect.is_valid or rect.is_empty:
                    continue
                rr = rect.round()
                if any(abs(rr.x0 - k.x0) < 2 and abs(rr.y0 - k.y0) < 2 and abs(rr.x1 - k.x1) < 2 and abs(rr.y1 - k.y1) < 2 for _pi, k in keep_rects if _pi == pi):
                    continue
                try:
                    text = page.get_textbox(rect).strip()
                except Exception:
                    continue
                if is_metadata_rect(page, rect, text):
                    try:
                        page.delete_annot(annot)
                        removed += 1
                    except Exception:
                        pass
    hl_final = os.path.join(out_dir, pn_x + "_highlight.pdf")
    doc.save(hl_final, garbage=4, deflate=True)
    doc.close()

    # PyMuPDF annot 序列化兼容: 二次 clean save, 避免 annot API 原生崩溃 (0xC0000005)
    d2 = fitz.open(hl_final)
    d2.save(hl_final + ".clean", garbage=4, deflate=True)
    d2.close()
    os.replace(hl_final + ".clean", hl_final)

    # --- 导出图片 (重新打开, 渲染不走 annot API) ---
    doc = fitz.open(hl_final)
    annot_pages = {}
    for pi in range(len(doc)):
        n = len(list(doc[pi].annots() or []))
        if n:
            annot_pages[pi + 1] = n
    for pi in annot_pages:
        pix = doc[pi - 1].get_pixmap(dpi=150)
        pix.save(os.path.join(out_dir, "%s_highlight_p%d.png" % (pn_x, pi)))
    pages_dir = os.path.join(out_dir, pn_x + "_highlight_pages")
    os.makedirs(pages_dir, exist_ok=True)
    for pi in range(len(doc)):
        pix = doc[pi].get_pixmap(dpi=100)
        pix.save(os.path.join(pages_dir, "page_%03d.jpg" % (pi + 1)))
    doc.close()
    os.remove(hl_tmp)
    return {
        "pn_x": pn_x,
        "text_ok": sum(1 for r in report if r[2].startswith("OK")),
        "annot_pages": annot_pages,
        "removed_9_rules": removed,
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--project-dir", default=None)
    parser.add_argument("--slide", type=int, default=None, help="只处理指定 slide (默认: 有 PDF 的全部 slide)")
    parser.add_argument("--only", default=None, help="只处理指定 Pn-x, 逗号分隔")
    parser.add_argument("--force", action="store_true", help="已存在输出也重跑")
    parser.add_argument("--no-render", action="store_true", help="跳过 slide 图导出")
    ns = parser.parse_args()
    root = project_root(ns)
    if not root or not os.path.isdir(root):
        print("错误: 项目根不存在 (--project-dir 或 TMA_PROJECT)")
        return 1
    pptx_path = None
    for f in os.listdir(root):
        if f.endswith(".pptx"):
            pptx_path = os.path.join(root, f)
            break
    if not pptx_path:
        print("错误: 项目根没有 PPTX")
        return 1
    pdf_dir = os.path.join(root, "_2_pdfs")
    out_base = os.path.join(root, "_highlight_nested")
    if not os.path.isdir(pdf_dir):
        print("错误: 无 _2_pdfs")
        return 1

    # Step 0: 导出 slide 图
    if not ns.no_render:
        renders = os.path.join(root, "_ppt_renders")
        n = render_ppt_slides(pptx_path, renders)
        print("[0] 导出 slide 图: %d 页 → %s" % (n, renders), flush=True)

    # vision report
    vpath = os.path.join(root, "_vision_report.json")
    vision = json.load(open(vpath, encoding="utf-8")) if os.path.exists(vpath) else None

    pdfs = sorted(f for f in os.listdir(pdf_dir) if f.endswith(".pdf") and re.match(r"P\d+-\d+\.pdf$", f))
    only = set(ns.only.split(",")) if ns.only else None
    slides = set()
    for f in pdfs:
        m = re.match(r"P(\d+)-(\d+)\.pdf", f)
        slides.add(int(m.group(1)))
    if ns.slide:
        slides = {ns.slide}

    results = []
    for slide in sorted(slides):
        slide_pdfs = [f for f in pdfs if f.startswith("P%d-" % slide)]
        if not slide_pdfs:
            continue
        print("\n[1] slide %d: 视觉提取 → %d 个 PDF" % (slide, len(slide_pdfs)), flush=True)
        visual = extract_slide_visual(pptx_path, slide, vision)
        print("    文本块 %d / 表格 %d / 图片形状 %d / 标号 %d" % (
            len(visual["text_blocks"]), len(visual["tables"]), len(visual["images"]), len(visual["citation_marks"])), flush=True)
        for pdf_file in slide_pdfs:
            if only and pdf_file.replace(".pdf", "") not in only:
                continue
            pn = pdf_file.replace(".pdf", "")
            hl_pdf = os.path.join(out_base, pn, pn + "_highlight.pdf")
            if not ns.force and os.path.exists(hl_pdf) and os.path.getsize(hl_pdf) > 10000:
                print("    [skip] %s (已有输出)" % pdf_file, flush=True)
                continue
            print("    [2] %s 应证 highlight..." % pdf_file, flush=True)
            r = highlight_one_pdf(os.path.join(pdf_dir, pdf_file), visual, os.path.join(out_base, pn))
            print("      文字 OK=%d / annot页=%s / 9铁律删除=%d" % (
                r["text_ok"], list(r["annot_pages"].keys()), r["removed_9_rules"]), flush=True)
            results.append(r)
    print("\n=== 完成: %d 个 Pn-x ===" % len(results), flush=True)
    with open(os.path.join(out_base, "_by_slide_meta.json"), "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    return 0


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    sys.exit(main())
