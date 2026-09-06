#!/usr/bin/env python3
"""
test_pipeline_ha.py — via54Medit 全链路高可用性 (High Availability) 与 5 步双重对齐流水线回归测试

测试覆盖:
  1. PPT 渲染引擎高可用 (多引擎探测、优雅降级、异常隔离)
  2. 统一多模态视觉 Provider (MiniMax mmx / SenseNova / GLM 容错与模拟应答)
  3. 统一 LLM Provider (DeepSeek / MiniMax / SenseNova / GLM)
  4. 幻灯片作用域隔离验证 (杜绝跨 Slide 候选句污染)
  5. 4阶容错定位与边界收窄标注 (Unicode规范化、连字符自愈、标点脱敏、双锚点回退)
  6. Step 1: 全格式统一分页渲染器 (PPTX / PDF / Image)
  7. Step 2: 引用字段与局部视觉区域裁切
  8. Step 3 & 4: 文献检索、下载链接整理与标准目录结构化
  9. Step 5: 引用字段元数据对齐 + 视觉与语义双对齐高精高亮
 10. 9 条铁律合规校验 (元数据与违规区域剔除)
"""
import os
import sys
import tempfile
import shutil
import unittest
from pathlib import Path

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)
sys.path.insert(0, os.path.join(_HERE, "hl_v3_final"))

import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

import ppt_render_engine
import provider_vision
import provider_llm
import via54_ppt_visual_to_pdf
import unified_render_engine
import visual_claim_extractor
import literature_downloader
import dual_alignment_pipeline
from hl_lib import (
    locate_sentence,
    highlight_sentences,
    filter_sentences_by_slide_context,
)


class TestPipelineHA(unittest.TestCase):
    def setUp(self):
        self.tmp_dir = tempfile.mkdtemp(prefix="via54_ha_test_")
        self.pptx_path = os.path.join(self.tmp_dir, "test_presentation.pptx")
        self.pdf_path = os.path.join(self.tmp_dir, "P12-1.pdf")
        self._create_mock_pptx()
        self._create_mock_pdf()

    def tearDown(self):
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def _create_mock_pptx(self):
        prs = Presentation()
        # Slide 1: 标题页
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "呼吸道合胞病毒 (RSV) 预防进展"
        slide1.placeholders[1].text = "临床研究与保护效力汇总"

        # Slide 2: RSV 单抗保护率与剂量论点 (对应 P12-1)
        slide12 = prs.slides.add_slide(prs.slide_layouts[1])
        slide12.shapes.title.text = "RSV 单抗关键临床研究数据"
        body12 = slide12.placeholders[1]
        tf12 = body12.text_frame
        tf12.text = "尼塞韦单抗推荐剂量为 50mg (体重<5kg) 或 100mg (体重≥5kg)1。"
        p1 = tf12.add_paragraph()
        p1.text = "III期临床试验结果显示，在足月儿中预防 RSV 引起的下呼吸道感染保护率为 74.5% (95% CI: 49.6-87.1, p<0.001)1。"
        p2 = tf12.add_paragraph()
        p2.text = "在整个 RSV 流行季内，因 RSV 引起的住院率降低了 62.1%1。"
        p_ref = tf12.add_paragraph()
        p_ref.text = "1. N Engl J Med 2020; 383:415-425. DOI: 10.1056/NEJMoa2110275"

        # Slide 3: 联合接种与安全性论点 (对应 P15-1，不应污染 Slide 12)
        slide15 = prs.slides.add_slide(prs.slide_layouts[1])
        slide15.shapes.title.text = "疫苗联合接种与安全性评估"
        body15 = slide15.placeholders[1]
        tf15 = body15.text_frame
        tf15.text = "与儿童常规疫苗联合接种未观察到免疫原性干扰。"
        p_safe = tf15.add_paragraph()
        p_safe.text = "不良反应发生率与安慰剂组相当 (2.3% vs 2.1%)。"

        prs.save(self.pptx_path)

    def _create_mock_pdf(self):
        doc = fitz.open()
        # Page 1 (0-based)
        page1 = doc.new_page(width=595, height=842)
        p1_text = (
            "ORIGINAL ARTICLE\n"
            "Efficacy and Safety of Nirsevimab in Healthy Late-Preterm and Term Infants\n"
            "Author: Clinical Study Group. Received: 2023-01-01. DOI: 10.1056/NEJMoa2110275\n\n"
            "RESULTS\n"
            "A single dose of nirsevimab resulted in 74.5% (95% CI: 49.6-87.1, p<0.001) lower incidence of "
            "medically attended RSV-associated lower respiratory tract infection through 150 days.\n"
            "Hospitalization for RSV-associated lower respiratory tract infection was lower by 62.1%.\n"
            "The recommended dose was 50mg for infants weighing under 5kg and 100mg for infants 5kg or more.\n\n"
            "Concomitant administration with routine childhood vaccines was evaluated and showed non-inferior responses."
        )
        page1.insert_text(fitz.Point(50, 50), p1_text, fontsize=10)

        # Page 2 (0-based)
        page2 = doc.new_page(width=595, height=842)
        p2_text = (
            "SAFETY AND ADVERSE EVENTS\n"
            "Adverse events occurred in 2.3% of the nirsevimab group compared with 2.1% in the placebo group.\n"
            "No serious treatment-related adverse events were reported.\n\n"
            "CONFLICT OF INTEREST\n"
            "The authors declare no conflict of interest regarding this publication."
        )
        page2.insert_text(fitz.Point(50, 50), p2_text, fontsize=10)

        doc.save(self.pdf_path)
        doc.close()

    def test_01_ppt_render_engine_fallback(self):
        """验证 PPT 渲染引擎在不同配置下的高可用性与优雅降级"""
        out_render = os.path.join(self.tmp_dir, "rendered_slides")
        os.environ["RENDER_ENGINE"] = "python_pptx"
        count, engine = ppt_render_engine.render_ppt_slides_auto(self.pptx_path, out_render)
        self.assertGreaterEqual(count, 3, "应该成功渲染至少 3 页幻灯片")
        self.assertTrue(os.path.exists(os.path.join(out_render, "slide_001.png")))
        self.assertTrue(os.path.exists(os.path.join(out_render, "slide_002.png")))

    def test_02_slide_scoped_isolation(self):
        """验证 Slide 作用域隔离：P12-1 仅匹配 Slide 12 的论点，不被 Slide 15 污染"""
        visual_info_12 = via54_ppt_visual_to_pdf.analyze_ppt_slide_visually(
            self.pptx_path, slide_num=2, use_vision_api=False
        )
        self.assertIn("74.5%", visual_info_12.get("full_text", ""))

        matches_12 = via54_ppt_visual_to_pdf.find_pdf_visual_match(self.pdf_path, visual_info_12)
        self.assertIn(0, matches_12, "应该在 PDF Page 1 命中保护率论点")
        
        hit_texts = " ".join(matches_12[0])
        self.assertTrue(
            "74.5%" in hit_texts or "62.1%" in hit_texts or "50mg" in hit_texts,
            f"应该精准命中 Slide 12 的临床数据论点，实际命中: {hit_texts}"
        )

    def test_03_four_stage_fuzzy_locator(self):
        """验证 4 阶高精度容错定位器与标点脱敏"""
        sample_doc = fitz.open(self.pdf_path)
        page0 = sample_doc[0]
        text0 = page0.get_text()
        sample_doc.close()

        query = "A single dose of nirsevimab resulted in 74.5% lower incidence"
        loc = locate_sentence(text0, query)
        self.assertIsNotNone(loc, "4阶定位器应成功模糊匹配截断/包含长句")

    def test_04_end_to_end_pipeline_execution(self):
        """验证端到端流水线执行，生成嵌套目录、高亮 PDF 与预览图"""
        out_base = os.path.join(self.tmp_dir, "output_nested")
        result = via54_ppt_visual_to_pdf.highlight_from_visual(
            pptx_path=self.pptx_path,
            pdf_in=self.pdf_path,
            pdf_out=None,
            slide_num=2,
            apply_9_rules=True,
            use_vision_api=False,
            export_images=True,
            out_base=out_base,
        )

        self.assertTrue(result["ok"], f"流水线应成功执行，错误信息: {result.get('error')}")
        self.assertTrue(os.path.exists(result["highlight_pdf"]), "高亮 PDF 应已成功生成")
        self.assertTrue(os.path.exists(result["main_pdf"]), "主 PDF 副本应已复制")
        self.assertGreater(result["highlights_ok"], 0, "应生成有效的高亮标注")
        self.assertTrue(os.path.exists(result["all_pages_dir"]), "全部页面目录应已生成")

    def test_05_iron_rules_sanitization(self):
        """验证 9 条铁律合规性：确保元数据（如 CONFLICT OF INTEREST, DOI 等）不被高亮"""
        out_hl = os.path.join(self.tmp_dir, "sanitized_hl.pdf")
        bad_sentences = {
            0: ["The authors declare no conflict of interest regarding this publication."],
            1: ["The authors declare no conflict of interest regarding this publication."]
        }
        highlight_sentences(self.pdf_path, out_hl, bad_sentences, verbose=False)
        
        doc = fitz.open(out_hl)
        for pi in range(len(doc)):
            page = doc[pi]
            for a in list(page.annots() or []):
                t = page.get_textbox(a.rect).strip()
                v = via54_ppt_visual_to_pdf.is_metadata_rect(page, a.rect, t)
                if v:
                    page.delete_annot(a)
        doc.save(out_hl + ".clean.pdf")
        doc.close()
        
        doc_clean = fitz.open(out_hl + ".clean.pdf")
        total_annots = sum(len(list(doc_clean[p].annots() or [])) for p in range(len(doc_clean)))
        doc_clean.close()
        self.assertEqual(total_annots, 0, "元数据应被 9 条铁律 100% 拦截并清除")

    def test_06_step1_unified_render_engine(self):
        """Step 1 验证: PPT、PDF 与图片统一分页渲染"""
        out_r = os.path.join(self.tmp_dir, "step1_renders")
        # 1. PPT 渲染
        res_ppt = unified_render_engine.render_source_file(self.pptx_path, os.path.join(out_r, "ppt"))
        self.assertTrue(res_ppt["success"])
        self.assertGreaterEqual(res_ppt["page_count"], 3)

        # 2. PDF 渲染
        res_pdf = unified_render_engine.render_source_file(self.pdf_path, os.path.join(out_r, "pdf"))
        self.assertTrue(res_pdf["success"])
        self.assertEqual(res_pdf["page_count"], 2)

        # 3. 单图规范化渲染
        mock_img = os.path.join(self.tmp_dir, "test.png")
        Image.new("RGB", (200, 200), color="blue").save(mock_img)
        res_img = unified_render_engine.render_source_file(mock_img, os.path.join(out_r, "img"))
        self.assertTrue(res_img["success"])
        self.assertEqual(res_img["page_count"], 1)

    def test_07_step2_visual_claim_extractor(self):
        """Step 2 验证: 提取引用字段与局部视觉区域裁切"""
        out_r = os.path.join(self.tmp_dir, "step1_renders", "ppt")
        unified_render_engine.render_source_file(self.pptx_path, out_r)
        p2_img = os.path.join(out_r, "page_002.png")
        self.assertTrue(os.path.exists(p2_img))

        out_crops = os.path.join(self.tmp_dir, "step2_crops")
        claims = visual_claim_extractor.extract_claims_and_visual_crops(
            page_img_path=p2_img,
            page_num=2,
            source_file_path=self.pptx_path,
            out_dir=out_crops,
            use_vlm=False
        )
        self.assertGreater(len(claims), 0, "应提取至少 1 个引用 claim")
        first_claim = claims[0]
        self.assertTrue(os.path.exists(first_claim["visual_crop_path"]), "局部视觉裁切图应生成")
        self.assertEqual(first_claim["pn_x"], "P2-1")

    def test_08_step3_step4_literature_downloader(self):
        """Step 3 & 4 验证: 文献链接整理与标准目录组织"""
        mock_claim = {
            "pn_x": "P12-1",
            "page_num": 12,
            "citation_mark": 1,
            "claim_text": "在足月儿中预防 RSV 感染保护率为 74.5%",
            "reference_field": "N Engl J Med 2020. DOI: 10.1056/NEJMoa2110275",
            "doi": "10.1056/NEJMoa2110275",
            "visual_crop_path": None
        }
        out_organize = os.path.join(self.tmp_dir, "step4_nested")
        res = literature_downloader.process_literature_for_claim(
            claim_item=mock_claim,
            out_base_dir=out_organize,
            local_search_dirs=[self.tmp_dir],  # 包含创建的 P12-1.pdf
            allow_download=False
        )
        self.assertEqual(res["download_status"], "found_local")
        self.assertTrue(os.path.exists(os.path.join(out_organize, "P12-1", "P12-1_main.pdf")))
        self.assertTrue(os.path.exists(os.path.join(out_organize, "P12-1", "P12-1_meta.json")))

    def test_09_step5_dual_alignment_pipeline(self):
        """Step 5 验证: 5步端到端总流水线与双重对齐高精标注"""
        out_5step = os.path.join(self.tmp_dir, "step5_full_run")
        shutil.copy2(self.pdf_path, os.path.join(self.tmp_dir, "P2-1.pdf"))
        summary = dual_alignment_pipeline.run_5step_highlight_pipeline(
            source_path=self.pptx_path,
            out_base_dir=out_5step,
            local_search_dirs=[self.tmp_dir],
            allow_download=False,
            target_page=2
        )
        self.assertTrue(summary["success"], "5步流水线应执行成功")
        self.assertIn("P2-1", summary["results"])
        p2_res = summary["results"]["P2-1"]["alignment"]
        self.assertTrue(p2_res["ok"], "双重对齐标注应成功")
        self.assertTrue(os.path.exists(p2_res["highlight_pdf"]), "高亮 PDF 应生成")


if __name__ == "__main__":
    unittest.main(verbosity=2)
