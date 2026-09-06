#!/usr/bin/env python3
"""
visual_claim_extractor.py — Step 2: 提取 a.引用字段 与 b.引用序号对应的视觉区域裁切图

输入:
  - 分页图片 (如 page_012.png)
  - 源文件 (PPTX / PDF / DOCX / 纯图片)
  - 页码 (1-based)

输出:
  - 引用字段集: 论点文本 (Claim)、引文标号 (1, 2)、底部参考文献字段 (Author, Journal, Year, DOI)
  - 视觉局部裁切图: {Pn-x}_claim_visual.png (精确从整页图片中裁切对应图表/数据点/标号区域)
  - 结构化元数据: {Pn-x}_claim.json
"""
import os
import re
import sys
import json
import shutil
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any
from PIL import Image

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)

from provider_vision import vision_analyze


def _extract_doi_from_text(text: str) -> Optional[str]:
    """从文本中提取标准 DOI"""
    if not text:
        return None
    m = re.search(r'10\.\d{4,9}/[^\s,;。\'"]+', text)
    return m.group(0).rstrip('.,;') if m else None


def _extract_pmid_from_text(text: str) -> Optional[str]:
    """从文本中提取 PMID"""
    if not text:
        return None
    m = re.search(r'PMID[:\s]*(\d{6,9})', text, re.IGNORECASE)
    return m.group(1) if m else None


def crop_image_bbox(image_path: str, bbox: List[float], out_crop_path: str, padding: int = 15) -> str:
    """
    根据 Bounding Box [x0, y0, x1, y1] 或 [ymin, xmin, ymax, xmax] 裁切图片。
    支持绝对像素坐标与 0~1000 / 0~1 归一化坐标。
    """
    os.makedirs(os.path.dirname(os.path.abspath(out_crop_path)), exist_ok=True)
    with Image.open(image_path) as im:
        w, h = im.size
        # 判定坐标格式
        c0, c1, c2, c3 = bbox[0], bbox[1], bbox[2], bbox[3]
        
        # 归一化 0~1
        if 0.0 <= c0 <= 1.0 and 0.0 <= c1 <= 1.0 and 0.0 <= c2 <= 1.0 and 0.0 <= c3 <= 1.0:
            x0, y0, x1, y1 = c0 * w, c1 * h, c2 * w, c3 * h
        # 归一化 0~1000 (部分 VLM 格式 [ymin, xmin, ymax, xmax])
        elif c0 <= 1000 and c1 <= 1000 and c2 <= 1000 and c3 <= 1000 and max(c0, c1, c2, c3) > 1.0 and min(w, h) > 1000:
            # ymin, xmin, ymax, xmax
            ymin, xmin, ymax, xmax = c0, c1, c2, c3
            x0, y0, x1, y1 = (xmin / 1000.0) * w, (ymin / 1000.0) * h, (xmax / 1000.0) * w, (ymax / 1000.0) * h
        else:
            x0, y0, x1, y1 = c0, c1, c2, c3

        # 确保 x0 < x1, y0 < y1
        min_x, max_x = min(x0, x1), max(x0, x1)
        min_y, max_y = min(y0, y1), max(y0, y1)

        # 加 padding 边距
        crop_box = (
            max(0, int(min_x - padding)),
            max(0, int(min_y - padding)),
            min(w, int(max_x + padding)),
            min(h, int(max_y + padding)),
        )

        # 保证有效宽高
        if crop_box[2] <= crop_box[0] or crop_box[3] <= crop_box[1]:
            crop_box = (0, 0, w, h)

        cropped = im.crop(crop_box)
        cropped.save(out_crop_path, "PNG")
    return out_crop_path


def extract_claims_and_visual_crops(
    page_img_path: str,
    page_num: int,
    source_file_path: Optional[str] = None,
    out_dir: Optional[str] = None,
    use_vlm: bool = True,
) -> List[Dict[str, Any]]:
    """
    Step 2 主流程:
    1. 提取所有引用标号 (1, 2, 3...)
    2. 提取 a-引用字段 (正文论点、统计数据、文献引文/DOI)
    3. 提取 b-引用区域裁切图 (从整页图片中精准裁切局部视觉图)
    """
    if not os.path.exists(page_img_path):
        return []

    target_out_dir = out_dir or os.path.dirname(os.path.abspath(page_img_path))
    os.makedirs(target_out_dir, exist_ok=True)

    extracted_items = []
    
    # 获取整图尺寸
    with Image.open(page_img_path) as im:
        img_w, img_h = im.size

    # 1. 尝试多模态大模型视觉理解 (VLM) 提取
    if use_vlm:
        prompt = """你是一位多模态临床文献比对专家。请仔细分析这张幻灯片/文档页面图片：
1. 找出页面中包含的所有文献引用标号（如上标数字 1, 2, [1], 1-3 等）。
2. 对于每一个引用标号，提取：
   - mark: 标号数字 (如 1)
   - claim_text: 标号所支撑的正文论点陈述（包括涉及的关键临床数据如有效率%、剂量mg、P值、生存期等）。
   - reference_field: 页面底部或脚注对应的完整参考文献条目（如作者、期刊名、年份、卷期、DOI等）。
   - visual_bbox: 包含该论点及相关图表/表格/数据的视觉区域坐标 [x0, y0, x1, y1] (像素坐标，或 0~1 归一化坐标)。

请输出 JSON 格式：
{
  "citations": [
    {
      "mark": 1,
      "claim_text": "在足月儿中预防 RSV 下呼吸道感染保护率为 74.5%...",
      "reference_field": "N Engl J Med 2020; 383:415-425. DOI: 10.1056/NEJMoa2110275",
      "visual_bbox": [50, 120, 700, 380]
    }
  ]
}"""
        try:
            resp = vision_analyze(page_img_path, prompt, json_mode=True, timeout=120)
            content = resp.get("content", "") if isinstance(resp, dict) else str(resp)
            json_m = re.search(r'\{.*\}', content, re.DOTALL)
            if json_m:
                data = json.loads(json_m.group(0))
                citations = data.get("citations", [])
                for idx, item in enumerate(citations):
                    mark = item.get("mark", idx + 1)
                    pn_x = f"P{page_num}-{mark}"
                    claim_text = item.get("claim_text", "").strip()
                    ref_field = item.get("reference_field", "").strip()
                    bbox = item.get("visual_bbox", [0, 0, img_w, img_h])
                    
                    crop_filename = f"{pn_x}_claim_visual.png"
                    crop_path = os.path.join(target_out_dir, crop_filename)
                    crop_image_bbox(page_img_path, bbox, crop_path)
                    
                    extracted_items.append({
                        "pn_x": pn_x,
                        "page_num": page_num,
                        "citation_mark": mark,
                        "claim_text": claim_text,
                        "reference_field": ref_field,
                        "doi": _extract_doi_from_text(ref_field) or _extract_doi_from_text(claim_text),
                        "pmid": _extract_pmid_from_text(ref_field),
                        "visual_crop_path": crop_path,
                        "bbox": bbox,
                        "source": "vlm"
                    })
        except Exception as e:
            print(f"  [visual_claim] VLM 视觉提取回退: {e}")

    # 2. 若 VLM 未提取到或源文件为 PPTX，使用结构化对象模型提取并补充
    if not extracted_items and source_file_path and source_file_path.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation
            prs = Presentation(source_file_path)
            if 1 <= page_num <= len(prs.slides):
                slide = prs.slides[page_num - 1]
                
                # 寻找脚注文献引用条目 (通常在 bottom 区域)
                footnote_texts = []
                body_claims = []
                
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    txt = shape.text_frame.text.strip()
                    if not txt:
                        continue
                        
                    top_pos = shape.top or 0
                    # 假定底部区域为参考文献字段
                    if top_pos > 5000000 or any(kw in txt.lower() for kw in ["1.", "doi:", "et al", "nejm", "lancet", "jco", "jama"]):
                        footnote_texts.append(txt)
                    else:
                        body_claims.append(txt)

                combined_ref = "\n".join(footnote_texts)
                combined_claim = "\n".join(body_claims)
                
                # 默认生成 Pn-1
                pn_x = f"P{page_num}-1"
                crop_filename = f"{pn_x}_claim_visual.png"
                crop_path = os.path.join(target_out_dir, crop_filename)
                
                # 默认裁切中间主体区域
                crop_bbox = [0, 0, img_w, int(img_h * 0.85)]
                crop_image_bbox(page_img_path, crop_bbox, crop_path)
                
                extracted_items.append({
                    "pn_x": pn_x,
                    "page_num": page_num,
                    "citation_mark": 1,
                    "claim_text": combined_claim,
                    "reference_field": combined_ref,
                    "doi": _extract_doi_from_text(combined_ref) or _extract_doi_from_text(combined_claim),
                    "pmid": _extract_pmid_from_text(combined_ref),
                    "visual_crop_path": crop_path,
                    "bbox": crop_bbox,
                    "source": "pptx_fallback"
                })
        except Exception as e:
            print(f"  [visual_claim] PPTX 结构化提取异常: {e}")

    # 3. 兜底保障: 若仍为空，生成默认整页裁切
    if not extracted_items:
        pn_x = f"P{page_num}-1"
        crop_path = os.path.join(target_out_dir, f"{pn_x}_claim_visual.png")
        shutil.copy2(page_img_path, crop_path)
        extracted_items.append({
            "pn_x": pn_x,
            "page_num": page_num,
            "citation_mark": 1,
            "claim_text": "",
            "reference_field": "",
            "doi": None,
            "pmid": None,
            "visual_crop_path": crop_path,
            "bbox": [0, 0, img_w, img_h],
            "source": "default"
        })

    # 保存各项元数据 json
    for item in extracted_items:
        json_path = os.path.join(target_out_dir, f"{item['pn_x']}_claim.json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(item, f, ensure_ascii=False, indent=2)

    return extracted_items


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 visual_claim_extractor.py <page_img> <page_num> [source_file] [out_dir]")
        sys.exit(1)
    p_img = sys.argv[1]
    p_num = int(sys.argv[2])
    s_file = sys.argv[3] if len(sys.argv) > 3 else None
    o_dir = sys.argv[4] if len(sys.argv) > 4 else None
    res = extract_claims_and_visual_crops(p_img, p_num, s_file, o_dir)
    print(f"Extracted {len(res)} claims & visual crops")
