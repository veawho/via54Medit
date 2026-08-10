#!/usr/bin/env python3
"""
ppt_expand.py — 6 步规则 #1 PPT 扩页工具 (2026-08-10)

按 6 步规则实现:
  1. 分析每页 shape 边界, 检测引用文献是否超出页面
  2. 自动扩页 (增加 slide 高度) 保证所有内容可见
  3. 分析引用文献文字颜色, 选合适底色 (保证对比度)
  4. 输出扩页后 PPT
  5. 可选: 导出每页 jpg/png

用法:
  python3.11 ppt_expand.py expand <input.pptx> <output.pptx> [--margin-pt 20]
  python3.11 ppt_expand.py audit <input.pptx>  # 只审计不修改
  python3.11 ppt_expand.py render <input.pptx> <out_dir> [--dpi 150]
"""
import os, sys, json, argparse, shutil
from pathlib import Path
from typing import Dict, List, Tuple, Optional

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ════════════════════════════════════════════════════════════════
# 常量
# ════════════════════════════════════════════════════════════════

# 6 步规则要求: 扩后底色必须保证引用文献内容可见
# 启发式: 如果引用文献文字色是浅色 (亮黄/白), 底色要深
# 如果是深色, 底色保持浅
LIGHT_TEXT_DARK_BG = (1, 1, 1)  # 白色文字
DARK_TEXT_LIGHT_BG = (1, 1, 1)  # 默认白色底

DEFAULT_MARGIN_PT = 20  # 扩页后边距 (pt)
DEFAULT_DPI = 150


# ════════════════════════════════════════════════════════════════
# 颜色辅助
# ════════════════════════════════════════════════════════════════

def _luminance(r: int, g: int, b: int) -> float:
    """估算颜色亮度 (0-1)"""
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255


def _analyze_text_colors(slide) -> List[Tuple[int, int, int]]:
    """提取 slide 上所有文字颜色"""
    colors = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        colors.append((rgb[0], rgb[1], rgb[2]))
                except (AttributeError, TypeError, ValueError):
                    pass
    return colors


def _is_dark_color(c: Tuple[int, int, int]) -> bool:
    """判断颜色是否深色"""
    return _luminance(*c) < 0.5


def _suggest_bg_color(text_colors: List[Tuple[int, int, int]]) -> Tuple[int, int, int]:
    """
    根据文字颜色建议底色 (保证对比度)
    - 多数文字是深色 → 浅底 (白)
    - 多数文字是浅色 → 深底 (黑)
    """
    if not text_colors:
        return (255, 255, 255)  # 默认白底
    dark_count = sum(1 for c in text_colors if _is_dark_color(c))
    if dark_count > len(text_colors) / 2:
        return (255, 255, 255)  # 多数深文字 → 白底
    else:
        return (32, 32, 32)  # 多数浅文字 → 深底


# ════════════════════════════════════════════════════════════════
# 边界检测
# ════════════════════════════════════════════════════════════════

def _get_slide_bounds(slide) -> Tuple[int, int]:
    """slide 宽高 (EMU)"""
    return slide.part.package.presentation_part.presentation.slide_width, \
           slide.part.package.presentation_part.presentation.slide_height


def _shape_bounds_emu(shape) -> Tuple[int, int, int, int]:
    """shape 边界 (left, top, right, bottom) in EMU"""
    return (shape.left or 0, shape.top or 0,
            (shape.left or 0) + (shape.width or 0),
            (shape.top or 0) + (shape.height or 0))


def _find_overflow_shapes(slide) -> List[Dict]:
    """
    找所有超出 slide 边界的 shape
    Returns: [{shape_idx, name, type, l, t, r, b, overflow: {top, bottom, left, right}}]
    """
    sw, sh = _get_slide_bounds(slide)
    overflows = []
    for i, shape in enumerate(slide.shapes):
        l, t, r, b = _shape_bounds_emu(shape)
        ov = {}
        if t < 0: ov["top"] = -t
        if b > sh: ov["bottom"] = b - sh
        if l < 0: ov["left"] = -l
        if r > sw: ov["right"] = r - sw
        if ov:
            overflows.append({
                "shape_idx": i,
                "name": shape.name,
                "type": str(shape.shape_type),
                "l": l, "t": t, "r": r, "b": b,
                "overflow": ov,
            })
    return overflows


# ════════════════════════════════════════════════════════════════
# 审计 (不修改)
# ════════════════════════════════════════════════════════════════

def audit_pptx(input_path: str) -> Dict:
    """
    审计 PPT 每页:
    - slide 宽高
    - 越界 shape 列表
    - 文字颜色 + 建议底色
    - 是否需要扩页
    """
    prs = Presentation(input_path)
    sw_emu, sh_emu = prs.slide_width, prs.slide_height
    sw_in, sh_in = sw_emu / 914400, sh_emu / 914400
    sw_pt, sh_pt = sw_emu / 12700, sh_emu / 12700

    report = {
        "input": input_path,
        "n_slides": len(prs.slides),
        "slide_size": {
            "width_in": round(sw_in, 2),
            "height_in": round(sh_in, 2),
            "width_pt": round(sw_pt, 1),
            "height_pt": round(sh_pt, 1),
        },
        "slides": [],
        "needs_expansion": 0,
    }

    for i, slide in enumerate(prs.slides):
        overflows = _find_overflow_shapes(slide)
        text_colors = _analyze_text_colors(slide)
        bg_suggestion = _suggest_bg_color(text_colors)

        slide_info = {
            "idx": i + 1,
            "n_shapes": len(slide.shapes),
            "n_overflow": len(overflows),
            "overflows": overflows,
            "text_colors_count": len(text_colors),
            "bg_suggestion": {
                "rgb": list(bg_suggestion),
                "hex": "#{:02x}{:02x}{:02x}".format(*bg_suggestion),
            },
        }
        if overflows:
            report["needs_expansion"] += 1
        report["slides"].append(slide_info)

    return report


# ════════════════════════════════════════════════════════════════
# 扩页 (修改)
# ════════════════════════════════════════════════════════════════

def expand_pptx(input_path: str, output_path: str, margin_pt: float = DEFAULT_MARGIN_PT) -> Dict:
    """
    自动扩页: 检测每页越界 shape, 增加 slide 高度保证内容可见
    """
    shutil.copy(input_path, output_path)
    prs = Presentation(output_path)
    sw_emu, sh_emu = prs.slide_width, prs.slide_height
    sw_pt, sh_pt = sw_emu / 12700, sh_emu / 12700
    margin_emu = int(margin_pt * 12700)

    log = {
        "input": input_path,
        "output": output_path,
        "original_size_pt": (round(sw_pt, 1), round(sh_pt, 1)),
        "slides_modified": 0,
        "slides_unchanged": 0,
        "details": [],
    }

    new_height_emu = sh_emu
    for i, slide in enumerate(prs.slides):
        overflows = _find_overflow_shapes(slide)
        if not overflows:
            log["slides_unchanged"] += 1
            log["details"].append({"slide": i + 1, "action": "unchanged"})
            continue

        # 算需要扩多少
        bottom_overflow = max((ov["overflow"].get("bottom", 0) for ov in overflows), default=0)
        needed = bottom_overflow + margin_emu
        if needed > 0:
            new_height_emu = max(new_height_emu, sh_emu + needed)
            log["details"].append({
                "slide": i + 1,
                "action": "expanded",
                "overflow_shapes": len(overflows),
                "expand_pt": round(needed / 12700, 1),
            })
            log["slides_modified"] += 1
        else:
            log["slides_unchanged"] += 1
            log["details"].append({"slide": i + 1, "action": "no_expansion_needed"})

    if new_height_emu > sh_emu:
        prs.slide_height = new_height_emu
        log["new_size_pt"] = (round(sw_pt, 1), round(new_height_emu / 12700, 1))
        prs.save(output_path)
    else:
        log["new_size_pt"] = log["original_size_pt"]
        # 即使没扩也保存一下 (确保是 pptx 格式)
        prs.save(output_path)

    return log


# ════════════════════════════════════════════════════════════════
# 渲染 (PPT → jpg)
# ════════════════════════════════════════════════════════════════

def render_pptx_images(input_path: str, out_dir: str, dpi: int = DEFAULT_DPI) -> List[str]:
    """
    把 PPT 每页渲染成 jpg
    需要 libreoffice 或 unoconv
    """
    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)

    # 尝试 libreoffice
    import subprocess, tempfile
    with tempfile.TemporaryDirectory() as tmp:
        # 转 PDF 先
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", tmp, input_path],
                check=True, capture_output=True, timeout=120
            )
        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
            print(f"  [render] libreoffice 失败: {e}")
            print("  提示: 安装 libreoffice (brew install --cask libreoffice)")
            return []

        pdf_path = os.path.join(tmp, os.path.basename(input_path).replace('.pptx', '.pdf'))
        if not os.path.isfile(pdf_path):
            print(f"  [render] PDF 未生成: {pdf_path}")
            return []

        # PDF → jpg
        import fitz
        import io
        from PIL import Image
        doc = fitz.open(pdf_path)
        out_files = []
        for i in range(len(doc)):
            try:
                pix = doc[i].get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
                img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                out = out_dir_p / f"slide_{i+1:03d}.jpg"
                img.save(out, "JPEG", quality=85)
                out_files.append(str(out))
            except Exception as e:
                print(f"  [render] page {i+1} fail: {e}")
        doc.close()
        return out_files


# ════════════════════════════════════════════════════════════════
# CLI
# ════════════════════════════════════════════════════════════════

def main():
    if not HAS_PPTX:
        print("需要 python-pptx: pip install python-pptx")
        sys.exit(1)

    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1]
    if cmd == "audit":
        if len(sys.argv) < 3:
            print("Usage: audit <input.pptx>")
            sys.exit(1)
        report = audit_pptx(sys.argv[2])
        print(json.dumps(report, ensure_ascii=False, indent=2))
    elif cmd == "expand":
        if len(sys.argv) < 4:
            print("Usage: expand <input.pptx> <output.pptx> [--margin-pt 20]")
            sys.exit(1)
        margin = DEFAULT_MARGIN_PT
        for i, a in enumerate(sys.argv):
            if a == "--margin-pt" and i + 1 < len(sys.argv):
                margin = float(sys.argv[i + 1])
        log = expand_pptx(sys.argv[2], sys.argv[3], margin)
        print(json.dumps(log, ensure_ascii=False, indent=2))
    elif cmd == "render":
        if len(sys.argv) < 4:
            print("Usage: render <input.pptx> <out_dir> [--dpi 150]")
            sys.exit(1)
        dpi = DEFAULT_DPI
        for i, a in enumerate(sys.argv):
            if a == "--dpi" and i + 1 < len(sys.argv):
                dpi = int(sys.argv[i + 1])
        files = render_pptx_images(sys.argv[2], sys.argv[3], dpi)
        print(f"渲染 {len(files)} 张")
        for f in files:
            print(f"  {f}")
    elif cmd == "all":
        # audit + expand + render 一次走完
        if len(sys.argv) < 4:
            print("Usage: all <input.pptx> <out_dir> [--margin-pt 20]")
            sys.exit(1)
        input_pptx = sys.argv[2]
        out_dir = sys.argv[3]
        os.makedirs(out_dir, exist_ok=True)
        out_pptx = os.path.join(out_dir, os.path.basename(input_pptx).replace('.pptx', '_expanded.pptx'))
        img_dir = os.path.join(out_dir, "images")

        print("=== Step 1: audit ===")
        report = audit_pptx(input_pptx)
        print(f"  {report['n_slides']} 页, {report['needs_expansion']} 页需扩")

        print("\n=== Step 2: expand ===")
        log = expand_pptx(input_pptx, out_pptx)
        print(f"  修改 {log['slides_modified']} 页, 不变 {log['slides_unchanged']} 页")
        print(f"  原: {log['original_size_pt']} -> 新: {log['new_size_pt']} (pt)")

        print("\n=== Step 3: render images ===")
        files = render_pptx_images(out_pptx, img_dir)
        print(f"  渲染 {len(files)} 张")
    else:
        print(f"未知命令: {cmd}")
        print(__doc__)
        sys.exit(1)


if __name__ == "__main__":
    main()
