#!/usr/bin/env python3
"""
ppt_render_engine.py — PPT → 图片 多引擎自动渲染 (部署新系统即用)

按优先级自动接入系统可用引擎:
  1. Microsoft PowerPoint COM   (ProgID: PowerPoint.Application)
  2. WPS 演示 COM                (ProgID: KWPP.Application, 接口兼容 PowerPoint)
  3. python-pptx + Pillow 近似渲染 (兜底, 任何平台可用)

Windows: COM 真实渲染 (保真度最高); 若系统有 PowerPoint/WPS 但缺 pywin32, 自动尝试 pip 安装。
macOS/Linux: 优先 soffice/libreoffice 真实渲染 (检测到即用, 全平台), 否则 python-pptx 近似渲染。
CJK 字体: 按平台探测 (Windows 微软雅黑 / macOS 苹方-简 / Linux Noto CJK), 保证中文近似渲染可读。

用法:
  from ppt_render_engine import render_ppt_slides_auto, detect_engines
  n, engine = render_ppt_slides_auto("D:/x.pptx", "D:/out")
"""
import os, io, sys, subprocess

# 引擎探测优先级
COM_ENGINES = [
    ("PowerPoint", "PowerPoint.Application"),
    ("WPS 演示", "KWPP.Application"),
]

# 渲染引擎偏好 (2026-09-04 用户规范): 默认 PowerPoint 并禁用其他引擎自动切换。
# 覆盖: 环境变量 RENDER_ENGINE=powerpoint|wps|libreoffice|python_pptx|auto
#   auto = 旧行为 (按可用性自动降级, 仅显式要求时启用)
def _engine_pref():
    return os.environ.get("RENDER_ENGINE", "powerpoint").strip().lower()

# 偏好引擎 → 引擎标识
_PREF_MAP = {
    "powerpoint": ("PowerPoint", "com", "PowerPoint.Application"),
    "ppt": ("PowerPoint", "com", "PowerPoint.Application"),
    "wps": ("WPS 演示", "com", "KWPP.Application"),
    "libreoffice": ("LibreOffice (soffice)", "soffice", None),
    "soffice": ("LibreOffice (soffice)", "soffice", None),
    "python_pptx": ("python-pptx 近似渲染", "python_pptx", ""),
    "python-pptx": ("python-pptx 近似渲染", "python_pptx", ""),
}


def _progid_available(progid):
    """注册表检查 COM ProgID 是否已注册 (Windows)"""
    if os.name != "nt":
        return False
    try:
        import winreg
        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, progid):
            return True
    except OSError:
        return False
    except Exception:
        return False


def _ensure_pywin32(progid_hint=""):
    """确保 win32com 可用; 检测到 COM 引擎但缺 pywin32 时自动安装"""
    try:
        import win32com.client  # noqa: F401
        return True
    except ImportError:
        pass
    if os.name != "nt":
        return False
    if progid_hint and not _progid_available(progid_hint):
        return False
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "pywin32"],
                       capture_output=True, timeout=300)
        import win32com.client  # noqa: F401
        return True
    except Exception:
        return False


def _com_probe(progid):
    """COM 试连 (会短暂启动应用), 成功返回 True"""
    try:
        import win32com.client
        app = win32com.client.DispatchEx(progid)  # DispatchEx: 强制新实例 (Dispatch 对 PowerPoint 有连接残留问题)
        try:
            app.Quit()
        except Exception:
            pass
        return True
    except Exception:
        return False


def _macos_powerpoint_available():
    """检查 macOS 系统中是否安装有 Microsoft PowerPoint"""
    if sys.platform != "darwin":
        return False
    try:
        r = subprocess.run(["osascript", "-e", 'id of app "Microsoft PowerPoint"'],
                           capture_output=True, text=True, timeout=5)
        return r.returncode == 0 and "com.microsoft.Powerpoint" in r.stdout
    except Exception:
        return False


def _find_soffice():
    """探测 LibreOffice 可执行文件 (soffice/libreoffice), 全平台"""
    import shutil
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    return None


def detect_engines():
    """返回可用引擎列表 [(name, kind, target), ...], kind: com | macos_ppt | soffice | python_pptx"""
    out = []
    if os.name == "nt":
        for name, progid in COM_ENGINES:
            if _progid_available(progid):
                if _ensure_pywin32(progid) and _com_probe(progid):
                    out.append((name, "com", progid))
    if sys.platform == "darwin" and _macos_powerpoint_available():
        out.append(("PowerPoint (macOS)", "macos_ppt", "com.microsoft.Powerpoint"))
    soffice = _find_soffice()
    if soffice:
        out.append(("LibreOffice (soffice)", "soffice", soffice))
    out.append(("python-pptx 近似渲染", "python_pptx", ""))
    return out


# ============ macOS 原生 PowerPoint 真实渲染 ============
def render_via_macos_powerpoint(pptx_path, out_dir, dpi=150):
    """macOS 下通过 AppleScript 控制原生 Microsoft PowerPoint 导出 PDF，再由 PyMuPDF 导出高清 PNG"""
    import tempfile
    import fitz
    abs_pptx = os.path.abspath(pptx_path)
    os.makedirs(out_dir, exist_ok=True)
    tmp_dir = tempfile.mkdtemp(prefix="ppt_mac_")
    tmp_pdf = os.path.join(tmp_dir, "slides.pdf")
    script = f'''
    tell application "Microsoft PowerPoint"
        set origApp to current application
        open POSIX file "{abs_pptx}"
        set thePres to active presentation
        save thePres in POSIX file "{tmp_pdf}" as save as PDF
        close thePres saving no
    end tell
    '''
    try:
        res = subprocess.run(["osascript", "-e", script], capture_output=True, text=True, timeout=300)
        if res.returncode != 0 or not os.path.exists(tmp_pdf):
            raise RuntimeError("PowerPoint AppleScript error: %s" % (res.stderr or res.stdout)[-200:])
        doc = fitz.open(tmp_pdf)
        n = 0
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=dpi)
            pix.save(os.path.join(out_dir, "slide_%03d.png" % i))
            n += 1
        doc.close()
        return n
    finally:
        import shutil
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ============ COM 真实渲染 ============
def render_via_com(progid, pptx_path, out_dir, width_px=1600):
    """用 PowerPoint/WPS COM 把每页 slide 导出为 PNG"""
    import win32com.client
    app = win32com.client.DispatchEx(progid)
    pres = None
    n = 0
    try:
        try:
            app.Visible = False
        except Exception:
            pass
        try:
            app.DisplayAlerts = False
        except Exception:
            pass
        pres = app.Presentations.Open(pptx_path, ReadOnly=True, Untitled=False, WithWindow=False)
        slides = pres.Slides
        total = slides.Count
        # 按页面比例求高度
        try:
            w = pres.PageSetup.SlideWidth
            h = pres.PageSetup.SlideHeight
        except Exception:
            w, h = 12192000, 6858000  # 16:9 EMU 兜底
        height_px = int(width_px * h / w) if w else int(width_px * 9 / 16)
        for i in range(1, total + 1):
            out_png = os.path.join(out_dir, "slide_%03d.png" % i)
            slides(i).Export(out_png, "PNG", width_px, height_px)
            n += 1
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass
    return n


# ============ LibreOffice 真实渲染 (全平台, 检测到即优先于近似) ============
def render_via_soffice(soffice, pptx_path, out_dir, dpi=120):
    """soffice --headless 转 PDF → PyMuPDF 渲染 PNG (真实排版, 含矢量/图表)"""
    os.makedirs(out_dir, exist_ok=True)
    import tempfile
    import fitz
    tmp = tempfile.mkdtemp(prefix="ppt_soffice_")
    pdf_path = os.path.join(tmp, "slides.pdf")
    try:
        r = subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                            "--outdir", tmp, pptx_path],
                           capture_output=True, text=True, timeout=300)
        if r.returncode != 0 or not os.path.exists(pdf_path):
            raise RuntimeError("soffice convert failed: %s" % (r.stderr or r.stdout)[-200:])
        doc = fitz.open(pdf_path)
        n = 0
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=dpi)
            pix.save(os.path.join(out_dir, "slide_%03d.png" % i))
            n += 1
        doc.close()
        return n
    finally:
        import shutil
        shutil.rmtree(tmp, ignore_errors=True)


# ============ python-pptx 近似渲染 (兜底) ============

# CJK 字体候选 (按平台): Windows 微软雅黑 / macOS 苹方·黑体 / Linux Noto·文泉驿
_FONT_CANDIDATES = [
    # Windows
    r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\msyhl.ttc", r"C:\Windows\Fonts\simhei.ttf",
    # macOS
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/System/Library/Fonts/Supplemental/Songti.ttc",
    # Linux (Noto CJK / 文泉驿)
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJKsc-Regular.otf",
    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
]


def _find_cjk_font():
    """按平台探测可用的 CJK 字体文件, 找不到返回 None (退化为内置位图字体)"""
    for cand in _FONT_CANDIDATES:
        if os.path.exists(cand):
            return cand
    # Linux 兜底: fc-match 查询
    if os.name != "nt":
        try:
            out = subprocess.run(["fc-match", "-f", "%{file}", "sans-serif:lang=zh"],
                                 capture_output=True, text=True, timeout=10)
            p = (out.stdout or "").strip()
            if p and os.path.exists(p):
                return p
        except Exception:
            pass
    return None


def render_via_python_pptx(pptx_path, out_dir, dpi=120):
    os.makedirs(out_dir, exist_ok=True)
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        print("  [warn] render 依赖缺失: %s (跳过 slide 图导出)" % e)
        return 0
    font_path = _find_cjk_font()
    prs = Presentation(pptx_path)
    EMU_IN = 914400.0
    n = 0
    for idx, slide in enumerate(prs.slides, start=1):
        w = int(prs.slide_width / EMU_IN * dpi)
        h = int(prs.slide_height / EMU_IN * dpi)
        img = Image.new("RGB", (max(w, 1), max(h, 1)), "white")
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    bio = io.BytesIO(shape.image.blob)
                    im = Image.open(bio).convert("RGB")
                    x0 = int(shape.left / EMU_IN * dpi); y0 = int(shape.top / EMU_IN * dpi)
                    x1 = int((shape.left + shape.width) / EMU_IN * dpi)
                    y1 = int((shape.top + shape.height) / EMU_IN * dpi)
                    if x1 > x0 and y1 > y0:
                        img.paste(im.resize((x1 - x0, y1 - y0)), (x0, y0))
                elif shape.has_table:
                    tbl = shape.table
                    x0 = int(shape.left / EMU_IN * dpi); y0 = int(shape.top / EMU_IN * dpi)
                    x1 = int((shape.left + shape.width) / EMU_IN * dpi)
                    y1 = int((shape.top + shape.height) / EMU_IN * dpi)
                    draw.rectangle([x0, y0, x1, y1], outline="black")
                    rows = len(tbl.rows); cols = len(tbl.columns)
                    rh = (y1 - y0) / rows if rows else 0
                    cw = (x1 - x0) / cols if cols else 0
                    for ri in range(rows):
                        for ci in range(cols):
                            cell = tbl.cell(ri, ci)
                            txt = (cell.text or "").strip()[:40]
                            cx0 = x0 + ci * cw; cy0 = y0 + ri * rh
                            draw.rectangle([cx0, cy0, cx0 + cw, cy0 + rh], outline="black")
                            if txt:
                                fnt = ImageFont.truetype(font_path, max(int(rh * 0.5), 8)) if font_path else ImageFont.load_default()
                                draw.text((cx0 + 2, cy0 + 2), txt, fill="black", font=fnt)
                elif shape.has_text_frame:
                    txt = (shape.text_frame.text or "").strip()
                    if not txt:
                        continue
                    x0 = int(shape.left / EMU_IN * dpi); y0 = int(shape.top / EMU_IN * dpi)
                    fnt = ImageFont.truetype(font_path, 12) if font_path else ImageFont.load_default()
                    draw.text((x0 + 2, y0 + 2), txt[:200], fill="black", font=fnt)
            except Exception:
                continue
        img.save(os.path.join(out_dir, "slide_%03d.png" % idx))
        n += 1
    return n


def _build_engine_list():
    """按引擎偏好构造待尝试引擎列表 (默认 powerpoint 自动适配 Windows COM 与 macOS 原生 PowerPoint)."""
    pref = _engine_pref()
    if pref == "auto":
        return detect_engines()
    spec = _PREF_MAP.get(pref)
    if spec is None:
        raise RuntimeError("未知 RENDER_ENGINE=%r (可选: powerpoint|wps|libreoffice|python_pptx|auto)" % pref)
    name, kind, progid = spec
    if kind == "com":
        if os.name == "nt":
            if not (_progid_available(progid) and _ensure_pywin32(progid) and _com_probe(progid)):
                raise RuntimeError("[render] 偏好引擎 %s (COM %s) 在 Windows 不可用" % (name, progid))
            return [(name, kind, progid)]
        elif sys.platform == "darwin" and _macos_powerpoint_available():
            return [("PowerPoint (macOS)", "macos_ppt", "com.microsoft.Powerpoint")]
        else:
            # 尝试检测是否有 LibreOffice 替代，否则提示
            soffice = _find_soffice()
            if soffice:
                return [("LibreOffice (soffice 备选)", "soffice", soffice)]
            return [("python-pptx 近似渲染 (兜底)", "python_pptx", "")]
    if kind == "soffice":
        p = _find_soffice()
        if not p:
            raise RuntimeError("[render] 偏好引擎 LibreOffice (soffice) 不可用")
        return [(name, kind, p)]
    return [(name, kind, "")]  # python_pptx 兜底恒可用


def render_ppt_slides_auto(pptx_path, out_dir, width_px=1600):
    """按引擎偏好渲染全部 slide, 返回 (count, engine_name).
    默认偏好 = PowerPoint (Windows COM / macOS 原生 PowerPoint).
    """
    os.makedirs(out_dir, exist_ok=True)
    try:
        engines = _build_engine_list()
    except Exception as e:
        print("  [render] %s" % str(e), flush=True)
        print("  [render] 提示: 引擎偏好由 RENDER_ENGINE 控制 (默认 powerpoint)", flush=True)
        return 0, "none"
    for name, kind, progid in engines:
        try:
            if kind == "com":
                print("  [render] 引擎=%s (COM %s)" % (name, progid), flush=True)
                n = render_via_com(progid, pptx_path, out_dir, width_px)
                if n > 0:
                    return n, name
            elif kind == "macos_ppt":
                print("  [render] 引擎=%s (AppleScript)" % name, flush=True)
                n = render_via_macos_powerpoint(pptx_path, out_dir)
                if n > 0:
                    return n, name
            elif kind == "soffice":
                print("  [render] 引擎=%s (%s)" % (name, progid), flush=True)
                n = render_via_soffice(progid, pptx_path, out_dir)
                if n > 0:
                    return n, name
            else:
                print("  [render] 引擎=%s (兜底)" % name, flush=True)
                n = render_via_python_pptx(pptx_path, out_dir)
                if n > 0:
                    return n, name
        except Exception as e:
            print("  [render] %s 失败: %s (尝试下一引擎)" % (name, str(e)[:100]), flush=True)
    return 0, "none"


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("out_dir")
    parser.add_argument("--width", type=int, default=1600)
    ns = parser.parse_args()
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    n, engine = render_ppt_slides_auto(ns.pptx, ns.out_dir, ns.width)
    print("渲染 %d 页, 引擎: %s" % (n, engine))
