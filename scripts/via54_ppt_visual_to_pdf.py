#!/usr/bin/env python3
"""
via54_ppt_visual_to_pdf.py — PPT视觉理解 → PDF语义应证 → Highlight 完整端到端流水线

核心升级:
  Step 1: 多引擎自动渲染 PPT (Microsoft PowerPoint COM / macOS 原生 / LibreOffice / python-pptx)
  Step 2: 统一多模态视觉大模型 (MiniMax mmx-cli / SenseNova / GLM) 深度解析 PPT 论点与图表数据
  Step 3: 泛化医学实体与统计数值提取 + 幻灯片作用域严格隔离 (消灭跨 Slide 候选句污染)
  Step 4: 基于 hl_lib v3 FINAL 4阶容错定位 + 像素级边界收窄 + 9 条铁律合规校验

输入: PPTX + PDF 路径 (如 P12-1.pdf)
输出: 规范化高亮 PDF ({Pn-x}_highlight.pdf) + 校验报告
"""
import os
import re
import sys
import json
import base64
import tempfile
import shutil
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)
sys.path.insert(0, os.path.join(_HERE, "hl_v3_final"))

import fitz

# 导入底层组件
from hl_lib import (
    highlight_sentences,
    locate_sentence,
    filter_sentences_by_slide_context,
)
from via54_highlight_v3_final import (
    is_metadata_rect,
    get_max_font_size,
    TITLE_FONT_SIZE,
    METADATA_ZONE_RATIO,
)
from ppt_render_engine import render_ppt_slides_auto
from provider_vision import vision_analyze, get_api_key, encode_image, get_image_mime
from provider_llm import llm_complete


def render_single_ppt_slide(pptx_path: str, slide_num: int, out_dir: Optional[str] = None) -> Optional[str]:
    """渲染指定 PPT 页码 (1-based) 为高清 PNG 图片"""
    if not os.path.exists(pptx_path):
        return None
        
    temp_dir = out_dir or tempfile.mkdtemp(prefix=f"ppt_slide_{slide_num}_")
    os.makedirs(temp_dir, exist_ok=True)
    
    # 检查是否已存在渲染结果
    expected_img = os.path.join(temp_dir, f"slide_{slide_num:03d}.png")
    if os.path.exists(expected_img):
        return expected_img
        
    count, engine = render_ppt_slides_auto(pptx_path, temp_dir)
    if count > 0 and os.path.exists(expected_img):
        return expected_img
        
    # 兜底查找
    cand_files = sorted(list(Path(temp_dir).glob(f"slide_*{slide_num:02d}*.png")) + list(Path(temp_dir).glob(f"*{slide_num}*.png")))
    if cand_files:
        return str(cand_files[0])
        
    return None


def analyze_ppt_slide_visually(pptx_path: str, slide_num: int, use_vision_api: bool = True) -> Dict[str, Any]:
    """
    Step 1: 视觉识别 PPT 第 N 页内容 (图表、数据点、中文论点、文献标号)
    """
    visual_info: Dict[str, Any] = {
        "slide_num": slide_num,
        "text_blocks": [],
        "citation_marks": [],
        "data_points": [],
        "full_text": "",
    }
    
    # 1. 尝试视觉大模型分析 (VLM)
    img_path = render_single_ppt_slide(pptx_path, slide_num) if use_vision_api else None
    if img_path and os.path.exists(img_path):
        prompt = """你是一位专业的医学多模态分析专家。请仔细观察这张 PPT 幻灯片图片，提取所有临床医学论点与图表数据：
1. text_blocks: 提取所有标题、结论要点、小结列表（保留完整语义，不要拆得过碎）。
2. citation_marks: 提取所有参考文献上标数字 (如 [1], 1, 2-3) 及其所支撑的具体论点句子 (context_text)。
3. data_points: 提取所有关键临床数据 (百分比 %、P值、OR/HR/RR、生存期月数、例数n/N、保护率、有效率、不良反应发生率等)。
4. tables_or_charts: 若有表格或统计图，简述其核心结论与对照组数据。

请以 JSON 格式输出：
{
  "text_blocks": [{"text": "论点句子", "type": "title|body|chart_summary"}],
  "citation_marks": [{"mark": 1, "context_text": "引用标号1所支撑的完整句子"}],
  "data_points": [{"text": "75.7%", "type": "percentage|p_value|ratio|number", "context": "包含该数据的完整论述"}]
}"""
        try:
            resp = vision_analyze(img_path, prompt, json_mode=True, timeout=120)
            content = resp.get("content", "") if isinstance(resp, dict) else str(resp)
            json_m = re.search(r'\{.*\}', content, re.DOTALL)
            if json_m:
                parsed = json.loads(json_m.group(0))
                visual_info["text_blocks"] = parsed.get("text_blocks", [])
                visual_info["citation_marks"] = parsed.get("citation_marks", [])
                visual_info["data_points"] = parsed.get("data_points", [])
        except Exception as e:
            print(f"  [提示] VLM 视觉解析返回异常: {e}, 启用 PPTX 结构化补充提取")

    # 2. 从 PPTX 对象模型提取结构化文本作为可靠补充
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        if 1 <= slide_num <= len(prs.slides):
            slide = prs.slides[slide_num - 1]
            slide_texts = []
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    if shape.has_table:
                        tbl_txts = []
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    tbl_txts.append(cell.text.strip())
                        if tbl_txts:
                            combined = " | ".join(tbl_txts)
                            slide_texts.append(combined)
                            visual_info["text_blocks"].append({"text": combined, "type": "table"})
                    continue
                    
                txt = shape.text_frame.text.strip()
                if not txt:
                    continue
                slide_texts.append(txt)
                
                # 补充 text_blocks
                if not any(tb.get("text") == txt for tb in visual_info["text_blocks"]):
                    visual_info["text_blocks"].append({"text": txt, "type": "body"})
                    
                # 补充数据点提取 (泛化通用医学数据)
                for dp_match in re.finditer(r'(\d+(?:\.\d+)?\s*(?:%|mg|ml|kg|月|年|days?|hours?|weeks?|HR|OR|RR|CI|p\s*[<>=]\s*\d+))', txt, re.IGNORECASE):
                    val = dp_match.group(0)
                    start = max(0, dp_match.start() - 30)
                    end = min(len(txt), dp_match.end() + 30)
                    visual_info["data_points"].append({
                        "text": val,
                        "type": "data",
                        "context": txt[start:end].strip()
                    })
                    
            visual_info["full_text"] = "\n".join(slide_texts)
    except Exception as e:
        print(f"  [警告] python-pptx 结构化解析失败: {e}")

    return visual_info


def find_pdf_visual_match(
    pdf_path: str,
    ppt_visual: Dict[str, Any],
    top_candidates_per_page: int = 5
) -> Dict[int, List[str]]:
    """
    Step 2: 在 PDF 中语义级搜索 PPT 视觉内容与论点
    完全摒弃硬编码专科词库，采用通用医学实体 + 临床数值 + 幻灯片上下文过滤
    """
    doc = fitz.open(pdf_path)
    slide_context = ppt_visual.get("full_text", "")
    if not slide_context:
        slide_context = " ".join(tb.get("text", "") for tb in ppt_visual.get("text_blocks", []))

    # 1. 搜集来自 Slide 的所有论点及数据点
    query_snippets = []
    for mark in ppt_visual.get("citation_marks", []):
        ctx = mark.get("context_text", "").strip()
        if len(ctx) >= 10:
            query_snippets.append(ctx)
            
    for dp in ppt_visual.get("data_points", []):
        ctx = dp.get("context", "").strip()
        if len(ctx) >= 10:
            query_snippets.append(ctx)
            
    for tb in ppt_visual.get("text_blocks", []):
        txt = tb.get("text", "").strip()
        if len(txt) >= 15:
            query_snippets.append(txt)
            
    query_snippets = list(dict.fromkeys(query_snippets))

    def normalize(s: str) -> str:
        return re.sub(r'\s+', ' ', s).strip().lower()

    sentences_map: Dict[int, List[str]] = {}

    for pi in range(len(doc)):
        page = doc[pi]
        page_text = page.get_text()
        if not page_text.strip():
            continue
            
        page_norm = normalize(page_text)
        raw_lines = [l.strip() for l in page_text.split('\n') if len(l.strip()) >= 15]

        # 策略 A: 完整长句直接包含匹配
        for qs in query_snippets:
            qs_norm = normalize(qs)
            if len(qs_norm) >= 25 and qs_norm in page_norm:
                sentences_map.setdefault(pi, []).append(qs)

        # 策略 B: 泛化医学词元与统计数值共现匹配 (取代旧版的 TA-TMA 硬编码)
        for line in raw_lines:
            # 过滤常见的文献元数据 (作者、机构、利益冲突、引用声明等)
            if any(meta in line.lower() for meta in ["conflict of interest", "author contributions", "received:", "accepted:", "doi.org/", "all rights reserved", "supplementary material"]):
                continue

            # 提取行内的有效数字与医学/科学词元
            line_tokens = set(re.findall(r'[a-zA-Z]{4,}|\d+(?:\.\d+)?%?', line.lower()))
            if not line_tokens:
                continue

            # 使用 slide_context 进行相关性判定
            slide_tokens = set(re.findall(r'[a-zA-Z]{4,}|\d+(?:\.\d+)?%?', slide_context.lower()))
            overlap = line_tokens.intersection(slide_tokens)
            
            # 如果数字和核心专业词汇同时命中
            numeric_hits = [t for t in overlap if re.match(r'^\d', t)]
            if len(overlap) >= 3 or (len(overlap) >= 2 and len(numeric_hits) >= 1):
                sentences_map.setdefault(pi, []).append(line)

    doc.close()

    # 策略 C: 对每页命中的候选句进行 Slide Context 深度打分与 Top-K 过滤
    filtered_map: Dict[int, List[str]] = {}
    for pi, cand_list in sentences_map.items():
        if not cand_list:
            continue
        dedup_cands = list(dict.fromkeys(cand_list))
        best_cands = filter_sentences_by_slide_context(
            candidates=dedup_cands,
            slide_context=slide_context,
            top_k=top_candidates_per_page
        )
        if best_cands:
            filtered_map[pi] = best_cands

    return filtered_map


def highlight_from_visual(
    pptx_path: str,
    pdf_in: str,
    pdf_out: Optional[str] = None,
    slide_num: Optional[int] = None,
    apply_9_rules: bool = True,
    use_vision_api: bool = True,
    export_images: bool = True,
    out_base: Optional[str] = None,
) -> Dict[str, Any]:
    """
    完整端到端流程: PPT视觉理解 → PDF精准语义定位 → v3 FINAL 4阶定位高亮 + 9 铁律合规校验
    """
    result: Dict[str, Any] = {
        "ok": False,
        "pn_x": "",
        "output_dir": "",
        "main_pdf": pdf_in,
        "highlight_pdf": "",
        "highlight_images": [],
        "all_pages_dir": "",
        "ppt_visual": {},
        "pdf_sentences": {},
        "highlights_ok": 0,
        "highlights_removed": 0,
        "violations": [],
    }

    if not os.path.exists(pptx_path):
        result["error"] = f"PPTX 文件不存在: {pptx_path}"
        return result
    if not os.path.exists(pdf_in):
        result["error"] = f"PDF 文件不存在: {pdf_in}"
        return result

    # 1. 规范化 Pn-x 命名与目标目录 (例如 P12-1)
    pdf_basename = os.path.basename(pdf_in).replace(".pdf", "")
    pn_x = pdf_basename
    m_old = re.match(r"Pn-S(\d+)_(\d+)$", pn_x)
    if m_old:
        pn_x = f"P{m_old.group(1)}-{m_old.group(2)}"
    m_std = re.match(r"P(\d+)-(\d+)", pn_x)
    
    # 自动解析 Slide 编号 (严格隔离幻灯片作用域)
    target_slide = slide_num
    if target_slide is None and m_std:
        target_slide = int(m_std.group(1))

    # 输出目录规划
    if pdf_out is None:
        pdf_in_dir = os.path.dirname(os.path.abspath(pdf_in))
        parent_dir = os.path.dirname(pdf_in_dir)
        if out_base is None:
            out_base = os.path.join(parent_dir, "_highlight_nested")
        pdf_out_dir = os.path.join(out_base, pn_x)
    else:
        pdf_out_dir = os.path.dirname(os.path.abspath(pdf_out)) or "."
        if pdf_out.endswith(".pdf"):
            pdf_out_dir = os.path.dirname(os.path.abspath(pdf_out)) or "."

    os.makedirs(pdf_out_dir, exist_ok=True)
    result["pn_x"] = pn_x
    result["output_dir"] = pdf_out_dir

    # 复制主 PDF 副本
    main_pdf_dest = os.path.join(pdf_out_dir, f"{pn_x}_main.pdf")
    if os.path.abspath(pdf_in) != os.path.abspath(main_pdf_dest):
        try:
            shutil.copy2(pdf_in, main_pdf_dest)
            result["main_pdf"] = main_pdf_dest
        except Exception:
            result["main_pdf"] = pdf_in
    else:
        result["main_pdf"] = pdf_in

    # 确定高亮 PDF 路径
    if pdf_out is None or not pdf_out.endswith(".pdf"):
        highlight_pdf = os.path.join(pdf_out_dir, f"{pn_x}_highlight.pdf")
    else:
        highlight_pdf = pdf_out
    result["highlight_pdf"] = highlight_pdf

    print(f"  [Pipeline] Pn-x: {pn_x} | Target Slide: {target_slide}")
    print(f"  [Pipeline] Highlight PDF Target: {highlight_pdf}")

    # 2. PPT 视觉理解 (严格限制在 target_slide, 杜绝全 PPT 广播污染)
    slides_to_process = [target_slide] if target_slide else [1]
    for sn in slides_to_process:
        print(f"  [Step 1] 分析 PPT Slide {sn} (VLM={use_vision_api})...")
        visual = analyze_ppt_slide_visually(pptx_path, sn, use_vision_api=use_vision_api)
        result["ppt_visual"][sn] = visual

    # 3. PDF 语义应证检索
    print(f"  [Step 2] 在 PDF 中进行跨语言与语义应证检索...")
    all_sentences: Dict[int, List[str]] = {}
    for sn, visual in result["ppt_visual"].items():
        page_sents = find_pdf_visual_match(pdf_in, visual)
        for pi, sents in page_sents.items():
            all_sentences.setdefault(pi, []).extend(sents)

    # 全局唯一去重
    deduped_map: Dict[int, List[str]] = {}
    seen = set()
    for pi, sents in all_sentences.items():
        for s in sents:
            key = s.strip()[:60]
            if key not in seen:
                seen.add(key)
                deduped_map.setdefault(pi, []).append(s)
    result["pdf_sentences"] = deduped_map

    total_sents = sum(len(v) for v in deduped_map.values())
    print(f"  [Step 2] 检索到 {total_sents} 个高置信度论点证据句，分布在 {len(deduped_map)} 页")

    # 4. v3 FINAL 4阶容错定位 + 像素级精准标注
    print(f"  [Step 3] 执行 v3 FINAL 逐行精确标注...")
    try:
        report = highlight_sentences(pdf_in, highlight_pdf, deduped_map, verbose=False)
        result["hl_lib_report"] = report
        result["highlights_ok"] = sum(1 for r in report if len(r) >= 3 and str(r[2]).startswith("OK"))
    except Exception as e:
        result["error"] = f"hl_lib 标注执行失败: {e}"
        return result

    # 5. 9 条铁律合规过滤
    if apply_9_rules and os.path.exists(highlight_pdf):
        print(f"  [Step 4] 执行 9 条铁律合规校验 (清理违规/元数据高亮)...")
        doc = fitz.open(highlight_pdf)
        for pi in range(len(doc)):
            page = doc[pi]
            annots = list(page.annots() or [])
            for annot in annots:
                if not annot.rect:
                    continue
                rect = annot.rect
                try:
                    text = page.get_textbox(rect).strip()
                except Exception:
                    continue

                violation = is_metadata_rect(page, rect, text)
                if violation:
                    try:
                        page.delete_annot(annot)
                        result["highlights_removed"] += 1
                        result["violations"].append((pi + 1, violation, text[:60]))
                    except Exception:
                        pass

        tmp_path = highlight_pdf + ".tmp"
        doc.save(tmp_path, garbage=4, deflate=True)
        doc.close()
        shutil.move(tmp_path, highlight_pdf)

    # 6. 导出预览图像
    if export_images and os.path.exists(highlight_pdf):
        print(f"  [Step 5] 导出高精度页面渲染图...")
        try:
            doc = fitz.open(highlight_pdf)
            highlight_pages = [pi for pi in range(len(doc)) if list(doc[pi].annots() or [])]
            
            for pi in highlight_pages:
                page = doc[pi]
                pix = page.get_pixmap(dpi=150)
                out_img = os.path.join(pdf_out_dir, f"{pn_x}_highlight_p{pi+1}.png")
                pix.save(out_img)
                result["highlight_images"].append(out_img)

            all_pages_dir = os.path.join(pdf_out_dir, f"{pn_x}_highlight_pages")
            os.makedirs(all_pages_dir, exist_ok=True)
            result["all_pages_dir"] = all_pages_dir
            for pi in range(len(doc)):
                page = doc[pi]
                pix = page.get_pixmap(dpi=100)
                out_img = os.path.join(all_pages_dir, f"page_{pi+1:03d}.jpg")
                pix.save(out_img)
            doc.close()
        except Exception as e:
            print(f"  [提示] 图像导出异常: {e}")

    result["ok"] = True
    return result


def main():
    import argparse
    parser = argparse.ArgumentParser(description="via54Medit PPT视觉理解与PDF高精标注流水线")
    parser.add_argument("pdf_in", help="输入 PDF (如 _2_pdfs/P12-1.pdf)")
    parser.add_argument("--pptx", help="PPTX 文件路径 (默认自动在父级目录查找)")
    parser.add_argument("--slide", type=int, default=None, help="指定 PPT Slide 编号")
    parser.add_argument("--no-vision", action="store_true", help="禁用 VLM 视觉大模型，仅用结构化提取")
    parser.add_argument("--no-rules", action="store_true", help="不应用 9 条铁律")
    parser.add_argument("--no-images", action="store_true", help="不导出页面预览图片")
    parser.add_argument("--out-dir", default=None, help="输出根目录")
    args = parser.parse_args()

    # 自动定位 PPT
    pptx_path = args.pptx
    if not pptx_path:
        for search_dir in [os.path.dirname(os.path.dirname(args.pdf_in)), os.path.dirname(args.pdf_in), "."]:
            if os.path.isdir(search_dir):
                for f in os.listdir(search_dir):
                    if f.endswith(".pptx") and not f.startswith("~$"):
                        pptx_path = os.path.join(search_dir, f)
                        break
            if pptx_path:
                break

    if not pptx_path:
        print("错误: 找不到关联的 PPTX 文件 (请使用 --pptx 指定)")
        return 1

    result = highlight_from_visual(
        pptx_path=pptx_path,
        pdf_in=args.pdf_in,
        pdf_out=None,
        slide_num=args.slide,
        apply_9_rules=not args.no_rules,
        use_vision_api=not args.no_vision,
        export_images=not args.no_images,
        out_base=args.out_dir,
    )

    print("\n" + "=" * 50)
    print("【via54Medit 视觉标注执行报告】")
    print(f"  成功状态: {result['ok']}")
    print(f"  标识编号: {result.get('pn_x', '')}")
    print(f"  输出目录: {result.get('output_dir', '')}")
    print(f"  高亮 PDF: {result.get('highlight_pdf', '')}")
    print(f"  有效标注数: {result['highlights_ok']}")
    print(f"  铁律拦截数: {result['highlights_removed']}")
    print("=" * 50)
    return 0 if result["ok"] else 1


if __name__ == "__main__":
    main()