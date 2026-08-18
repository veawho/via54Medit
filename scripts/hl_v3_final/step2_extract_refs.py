#!/usr/bin/env python3
"""Step 2: 从 PPTX 提取每个 slide 的文献引用字段 → Pn-x 映射
用法: python3 step2_extract_refs.py <ppt_path> <out_json>
输出: [{slide, num, text}] 与现有 _citation_table/tma_citation_table.json 同构
规则: 引用片段 = "N. 文本" 且文本含年份(19xx/20xx)或 "et al"/中文期刊特征;
      支持同行多引用(如 slide23 整行 1..26)与跨段落"""
import sys, json, re
from pptx import Presentation

_REF_SEG = re.compile(r'(19|20)\d{2}')
# 编号: "N." 后必须跟大写字母或中文(排除 6.2/43.6 等数字)
_NUM = re.compile(r'(\d{1,2})\.\s*(?=[A-Z\u4e00-\u9fff])')

def is_ref(seg):
    """引用片段特征: 含年份 或 et al/等 或 中文长文本(无年份指南)"""
    if _REF_SEG.search(seg):
        return True
    if 'et al' in seg or '等' in seg:
        return True
    # 中文引用(期刊/指南, 可能无年份): 中文长度 >= 10 (排除正文短列表如 "突变，基因多态性")
    zh = len(re.findall(r'[\u4e00-\u9fff]', seg))
    return zh >= 10

def split_refs(text):
    """把含多个 "N. 引用" 的文本拆成 [(num, text), ...]"""
    nums = list(_NUM.finditer(text))
    out = []
    for i, m in enumerate(nums):
        start = m.end()
        end = nums[i + 1].start() if i + 1 < len(nums) else len(text)
        seg = text[start:end].strip()
        if not is_ref(seg):
            continue
        num = int(m.group(1))
        # 递归拆分: 片段内部若仍有 "N. 大写/中文" 模式(如 "6.2. Zheng" 合并), 继续拆
        sub = split_refs(seg)
        if sub and any(True for _ in sub):
            # 子拆分成功且首个子项起点 < 片段长度的一半, 说明确实合并了
            first_start = _NUM.search(seg).start() if _NUM.search(seg) else 0
            if sub[0][0] != num and first_start < len(seg) // 2:
                out.extend(sub)
                continue
        out.append((num, seg))
    # 去重(同 num 同文本)
    seen = set()
    dedup = []
    for num, seg in out:
        k = (num, re.sub(r'[\s\u3000]+', '', seg))
        if k not in seen:
            seen.add(k)
            dedup.append((num, seg))
    return dedup

def extract_refs(ppt_path):
    prs = Presentation(ppt_path)
    refs = []
    for si, slide in enumerate(prs.slides, start=1):
        seen = set()
        for shape in slide.shapes:
            texts = []
            if shape.has_text_frame:
                # 段落级拆分(段落边界处编号可跨行续接)
                for para in shape.text_frame.paragraphs:
                    t = ''.join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
            # 先按段落各自拆分, 段落间编号断续时合并处理
            joined = '\n'.join(texts)
            for num, seg in split_refs(joined):
                norm_seg = re.sub(r'[\s\u3000]+', '', seg)
                # 同 slide 同 num: 互相包含则保留更长者, 否则(如 slide31 的 Laurence/Jiang)都保留
                dup = [r for r in refs if r['slide'] == si and r['num'] == num]
                merged = False
                for r in dup:
                    rn = re.sub(r'[\s\u3000]+', '', r['text'])
                    if norm_seg in rn or rn in norm_seg:
                        if len(norm_seg) > len(rn):
                            r['text'] = seg
                        merged = True
                        break
                if not merged:
                    refs.append({'slide': si, 'num': num, 'text': seg})
    refs.sort(key=lambda r: (r['slide'], r['num']))
    return refs

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('usage: step2_extract_refs.py <ppt> <out_json>')
        sys.exit(1)
    refs = extract_refs(sys.argv[1])
    with open(sys.argv[2], 'w', encoding='utf-8') as f:
        json.dump(refs, f, ensure_ascii=False, indent=2)
    print(f'extracted {len(refs)} refs → {sys.argv[2]}')
