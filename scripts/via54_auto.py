#!/usr/bin/env python3
"""
via54_auto.py — 自然语言一键全自动管线 (任何设备部署即用)

自然语言示例:
  "帮我识别 D:/x/某PPT.pptx 中的文献引用，下载文献，并进行highlight"
  "下载并高亮 C:/文献/方案.pptx 的参考文献"

流程 (编排 deps_auto + ppt_render_engine + 下载器 + by-slide highlight):
  [0] 环境自检+自动安装依赖 (deps_auto.ensure_env)
  [1] 渲染全部 PPT slide 图 (_ppt_renders/, 自动接入 PowerPoint/WPS COM)
  [2] 提取文献引用 (逐 slide 视觉提取 → _refs.json + 8列 CSV 草稿)
  [3] 下载 PDF (级联 OA + CrossRef 重解析 + Sci-Hub; **1 小时硬限**, 超时剩余保留链接)
  [4] 整理下载目录 (_literature_citation_index/Pn-x/Pn-x_main.pdf)
  [5] 逐 slide 视觉分析 → highlight plan (_highlight_plans/slide_NNN_plan.json)
  [6] 按 slide 顺序完成 highlight (_highlight_nested/Pn-x/ 四类应证)
  [7] 交付报告 (_TMA_highlight_交付报告.md + _citations_8col.csv + _人工下载清单.md)

用法:
  python via54_auto.py "帮我识别 X.pptx 中的文献引用，下载文献，并进行highlight"
  python via54_auto.py --ppt X.pptx --download --highlight
  python via54_auto.py "自然语言" --budget 3600 --skip-highlight
"""
import os, re, sys, json, time, shutil, subprocess, argparse

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

DEFAULT_BUDGET_S = 3600  # 下载总时间硬限: 1 小时

# ============ 自然语言解析 ============
# 优先绝对路径 (盘符:/ 开头, 允许中文目录), 其次相对路径/裸文件名
# 路径字符排除中文标点/空白, 避免自然语言前缀被吞
PPTX_RE = re.compile(
    r"([A-Za-z]:[\\/][^\s，,。\"']*?\.pptx|"   # 盘符绝对路径
    r"[\\/][^\s，,。\"']*?\.pptx|"              # / 或 \ 开头
    r"[A-Za-z0-9_\-\.\\/]+?\.pptx)", re.I)


def parse_nl(text):
    """从自然语言提取 PPT 路径 + 意图步骤"""
    low = text.lower()
    m = PPTX_RE.search(text)
    pptx = m.group(1).strip().strip("，,。") if m else None
    return {
        "pptx": pptx,
        "render": True,  # 恒渲染 (渲染 PPT 图片是流程前提)
        "extract": True,
        "download": ("下载" in text or "download" in low or "downloaded" in low),
        "highlight": ("highlight" in low or "高亮" in text or "标注" in text),
        "raw": text[:200],
    }


# ============ 项目根解析 ============
def resolve_project(pptx, project_dir):
    if project_dir:
        return os.path.abspath(project_dir)
    if pptx:
        return os.path.dirname(os.path.abspath(pptx))
    return os.getcwd()


# ============ Step 1: 渲染 ============
def step_render(pptx, root):
    from ppt_render_engine import render_ppt_slides_auto
    renders = os.path.join(root, "_ppt_renders")
    n, engine = render_ppt_slides_auto(pptx, renders, width_px=1600)
    print("[1] 渲染 PPT 图: %d 页, 引擎=%s" % (n, engine), flush=True)
    return n


# ============ Step 2: 提取引用 ============
def _iter_shapes(shapes):
    for sh in shapes:
        if getattr(sh, 'shape_type', None) == 6:
            yield from _iter_shapes(sh.shapes)
        else:
            yield sh


def _sup_marks(para):
    """上标 run → 引用标号 (支持复合 "4,6"/"1-3"; 排除 C3/C5 前导字母)"""
    out = []
    full = para.text or ""
    for run in para.runs:
        txt = (run.text or "").strip()
        try:
            rPr = run.font._rPr
        except Exception:
            rPr = None
        if rPr is None or rPr.get("baseline") != "30000":
            continue
        nums = []
        for part in re.split(r"[,，;；]", txt):
            part = part.strip()
            if not part:
                continue
            if part.isdigit() and 1 <= int(part) <= 200:
                nums.append(int(part))
            elif "-" in part:
                try:
                    a, b = part.split("-")[:2]
                    if a.isdigit() and b.isdigit():
                        lo, hi = int(a), min(int(b), 200)
                        nums.extend(range(lo, hi + 1))
                except Exception:
                    pass
        for num in nums:
            pos = 0
            while True:
                idx = full.find(txt, pos)
                if idx < 0:
                    break
                prev = full[idx - 1] if idx > 0 else ""
                if prev and (prev.isalpha() or prev.isdigit()):
                    pos = idx + len(txt)
                    continue
                out.append((num, full[:120]))
                break
    return out



def _marks_in_para(para):
    marks = _sup_marks(para)
    full = para.text or ''
    for m in re.finditer(r'([\u4e00-\u9fff]{2,6})([1-9][0-9]?)(?=[,，;；]|$)', full):
        marks.append((int(m.group(2)), full[:120]))
    # 注: 括号 [n]/(n) 形式易误抓列表序号, 不采用; 引用标号以上标/中文+数字为准
    return marks


REF_LIST_PAT = re.compile(
    r"(?:^|[.\n])\s*(\d{1,3})\.\s+"
    r"([A-Za-z\u4e00-\u9fff][^。]{5,240}?(?:19|20)\d{2}[^。]{0,120}?)"
    r"(?=(?:\s+\d{1,3}\.\s+[A-Za-z\u4e00-\u9fff])|$)"
)


def extract_ref_list(pptx_path):
    """参考文献列表页解析: 扫描全部 slide 文本, 提取 'N. 完整引文' 编号列表"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    refs = {}
    for slide in prs.slides:
        for sh in _iter_shapes(slide.shapes):
            try:
                if sh.has_text_frame:
                    txt = sh.text_frame.text or ''
                elif sh.has_table:
                    txt = ' '.join((c.text or '') for r in sh.table.rows for c in r.cells)
                else:
                    continue
            except Exception:
                continue
            if not re.search(r'\d{1,3}\.\s+[A-Za-z]', txt):
                continue
            for m in REF_LIST_PAT.finditer(txt):
                num = int(m.group(1))
                cit = re.sub(r'\s+', ' ', m.group(2)).strip()
                if 1 <= num <= 500 and len(cit) >= 15:
                    refs.setdefault(num, cit)
    return refs


def deep_extract_refs(pptx_path):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    refs = {}
    for si, slide in enumerate(prs.slides, start=1):
        for sh in _iter_shapes(slide.shapes):
            try:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        for num, ctx in _marks_in_para(para):
                            refs.setdefault('P%d-%d' % (si, num), ctx)
                if sh.has_table:
                    for r in sh.table.rows:
                        for c in r.cells:
                            if not c.has_text_frame:
                                continue
                            for para in c.text_frame.paragraphs:
                                for num, ctx in _marks_in_para(para):
                                    refs.setdefault('P%d-%d' % (si, num), ctx)
            except Exception:
                continue
    return refs


def step_extract_refs(pptx, root):
    """深度提取引用 (上标/中文+数字/括号标号) + 融合 vision report 标号"""
    refs = deep_extract_refs(pptx)
    vpath = os.path.join(root, '_vision_report.json')
    if os.path.exists(vpath):
        vision = json.load(open(vpath, encoding='utf-8'))
        for sk, sd in (vision.get('slides') or {}).items():
            try:
                slide_num = int(sd.get('slide_num', sk))
            except Exception:
                continue
            for mid, mark in (sd.get('citation_marks') or {}).items():
                ctx = (mark.get('context') or '').strip()[:120]
                for part in re.split(r'[,，]', str(mid)):
                    part = part.strip()
                    nums = []
                    if '-' in part:
                        a, b = part.split('-')[:2]
                        try:
                            nums = list(range(int(a), int(b) + 1))
                        except Exception:
                            nums = []
                    else:
                        try:
                            nums = [int(part)]
                        except Exception:
                            nums = []
                    for num in nums:
                        key = 'P%d-%d' % (slide_num, num)
                        if ctx:
                            refs[key] = ctx
                        else:
                            refs.setdefault(key, '')
    # 参考文献列表 → 完整引文字典 (下载/报告用)
    full_refs = extract_ref_list(pptx)
    full_path = os.path.join(root, "_references_full.json")
    json.dump(full_refs, open(full_path, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    refs_path = os.path.join(root, "_refs.json")
    json.dump(refs, open(refs_path, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    print("[2] 提取引用: %d 条 (标号) + %d 条完整引文 → %s" % (len(refs), len(full_refs), root), flush=True)
    return refs

def _safe_remove(path):
    try:
        if os.path.exists(path):
            os.remove(path)
    except Exception:
        pass


def _verify_download(ref, citation, pdf_path):
    """下载后内容核验: 期刊/年份/作者 (复用 tma_verify_pdfs.check), mismatch 即删"""
    try:
        import tma_verify_pdfs as vp
        r = vp.check(ref, citation, pdf_path)
        return r.get("verdict") != "mismatch"
    except Exception:
        return True  # 核验异常时保守保留


def _pubmed_term_download(citation, out_path):
    """正文句含英文医学术语时, 用 PubMed ESearch 定位候选 PMID → EuropePMC 下载

    返回 True 若下载成功且为有效 PDF (内容核验由调用方执行)"""
    import urllib.parse as _up
    terms = re.findall(r"\b[A-Z]{2,12}\b", citation or "")
    terms = [t for t in terms if t not in ("TMA", "PPT") and len(t) >= 3][:4]
    if len(terms) < 2:
        return False
    q = " AND ".join(terms)
    try:
        d = json.loads(__import__("urllib.request").urlopen(
            "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi?db=pubmed&retmode=json&retmax=3&term="
            + _up.quote(q), timeout=25).read().decode("utf-8", "replace"))
        ids = d.get("esearchresult", {}).get("idlist", [])
        if not ids:
            return False
        for pmid in ids[:2]:
            ep = json.loads(__import__("urllib.request").urlopen(
                "https://www.ebi.ac.uk/europepmc/webservices/rest/search?query=EXT_ID:%s&resultType=core&format=json" % pmid,
                timeout=25).read().decode("utf-8", "replace"))
            res = (ep.get("resultList") or {}).get("result") or []
            for r in res:
                if r.get("isOpenAccess") == "Y" and r.get("pmcid"):
                    url = "https://europepmc.org/articles/%s/pdf/main.pdf" % r["pmcid"]
                    try:
                        req = __import__("urllib.request").Request(url, headers={
                            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"})
                        data = __import__("urllib.request").urlopen(req, timeout=90).read()
                        if data[:4] == b"%PDF" and len(data) > 5000:
                            open(out_path, "wb").write(data)
                            return True
                    except Exception:
                        continue
    except Exception:
        return False
    return False


def _assoc_full(full_cit, context, pdf_path):
    """双核验关联: ①完整引文自洽 (期刊/年份/作者) ②context 英文术语出现在 PDF 前 3 页"""
    try:
        import tma_verify_pdfs as vp
        r = vp.check("ref", full_cit, pdf_path)
        if r.get("verdict") == "mismatch":
            return False
    except Exception:
        pass
    terms = re.findall(r"\b[A-Z]{2,12}\b", context or "")
    terms = [t for t in terms if len(t) >= 3 and t not in ("TMA", "PPT", "HSCT", "HCT", "CI", "OR", "HR")][:5]
    if not terms:
        return True
    try:
        import fitz
        doc = fitz.open(pdf_path)
        txt = " ".join(doc[i].get_text() for i in range(min(3, len(doc)))).lower()
        doc.close()
        return sum(1 for t in terms if t.lower() in txt) >= 1
    except Exception:
        return True


# ============ Step 3: 下载 (1 小时硬限) ============
def step_download(root, budget_s=DEFAULT_BUDGET_S, limit=0):
    """级联 + 重解析下载, 超预算即停, 剩余保留链接"""
    os.environ["TMA_PROJECT"] = root
    import tma_cascade_download as cd
    import tma_download_round2 as rd

    refs_path = os.path.join(root, "_refs.json")
    refs = json.load(open(refs_path, encoding="utf-8"))
    # 注: 参考文献列表为全局编号, 而 P{slide}-{num} 的 num 是页内标号 (PPT 每页重新编号),
    #     两者无结构映射, 不做自动替换, 避免下载错配; 完整引文保留于 _references_full.json 供报告/人工核对。
    doi_map_path = os.path.join(root, "_doi_map.json")
    doi_map = json.load(open(doi_map_path, encoding="utf-8")) if os.path.exists(doi_map_path) else {}
    pdf_dir = os.path.join(root, "_2_pdfs")
    os.makedirs(pdf_dir, exist_ok=True)

    missing = []
    for k in sorted(refs, key=lambda x: (int(x.split("-")[0][1:]), int(x.split("-")[1]))):
        p = os.path.join(pdf_dir, k + ".pdf")
        if not (os.path.exists(p) and os.path.getsize(p) > 5000):
            missing.append(k)
    if limit:
        missing = missing[:limit]
    print("[3] 待下载: %d 条 (预算 %d s)" % (len(missing), budget_s), flush=True)

    start = time.monotonic()
    results = []
    budget_hit = False

    # --- 全文库优先: 参考文献列表完整引文 (准确字段, 高成功率), 用一半预算 ---
    full_path = os.path.join(root, "_references_full.json")
    full_ok = 0
    if os.path.exists(full_path):
        full_refs = json.load(open(full_path, encoding="utf-8"))
        full_budget = budget_s // 2
        print("  [全文库] 参考文献列表 %d 条完整引文直连下载 (预算 %d s)..." % (len(full_refs), full_budget), flush=True)
        for num, cit in full_refs.items():
            if time.monotonic() - start > full_budget:
                budget_hit = True
                print("    ⏰ 全文库预算耗尽, 剩余保留链接人工下载", flush=True)
                break
            refn = "ref%s" % num
            out_path = os.path.join(pdf_dir, refn + ".pdf")
            if os.path.exists(out_path) and os.path.getsize(out_path) > 5000:
                continue
            print("    [%s] %s..." % (refn, str(cit)[:50]), flush=True)
            try:
                r1 = cd.process_ref(refn, str(cit), {}, pdf_dir, 0.3)
                if r1.get("status") == "ok" or (os.path.exists(out_path) and os.path.getsize(out_path) > 5000):
                    print("      -> OK %s" % r1.get("source", "exists"), flush=True)
                    full_ok += 1
                else:
                    try:
                        r2 = rd.process_ref(refn, str(cit), {}, pdf_dir, 0.3)
                        if r2.get("status") == "ok" or (os.path.exists(out_path) and os.path.getsize(out_path) > 5000):
                            print("      -> OK %s" % r2.get("source", "exists"), flush=True)
                            full_ok += 1
                        else:
                            print("      -> FAILED (保留链接人工下载)", flush=True)
                    except Exception as e:
                        print("      -> round2 异常: %s" % str(e)[:60], flush=True)
            except Exception as e:
                print("      -> 异常: %s" % str(e)[:60], flush=True)

    # --- 标号 refs 下载 (剩余预算) ---
    for i, ref in enumerate(missing):
        if time.monotonic() - start > budget_s:
            budget_hit = True
            print("    ⏰ 预算耗尽, 停止下载, 剩余 %d 条保留链接人工下载" % (len(missing) - i), flush=True)
            break
        citation = refs[ref]
        print("  [%d/%d] %s: %s..." % (i + 1, len(missing), ref, citation[:40]), flush=True)
        out_path = os.path.join(pdf_dir, ref + ".pdf")
        ok = False
        # 编号命中完整引文时, 优先关联下载 (双核验防错配)
        try:
            _num = int(ref.split("-")[1])
        except Exception:
            _num = None
        if _num and str(_num) in full_refs:
            _full_cit = str(full_refs[str(_num)])
            try:
                r0 = cd.process_ref(ref, _full_cit, {}, pdf_dir, 0.3)
                if r0.get("status") == "ok" or (os.path.exists(out_path) and os.path.getsize(out_path) > 5000):
                    if _assoc_full(_full_cit, citation, out_path):
                        print("    -> OK 完整引文关联 %s" % r0.get("source", "exists"), flush=True)
                        results.append({"ref": ref, "status": "ok", "source": "full_assoc:%s" % r0.get("source")})
                        ok = True
                    else:
                        _safe_remove(out_path)
                        print("    -> 双核验失败, 回退 context 下载", flush=True)
                else:
                    print("    -> 完整引文下载失败, 回退 context", flush=True)
            except Exception as e:
                print("    -> 完整引文关联异常: %s" % str(e)[:60], flush=True)
        if not ok:
            try:
                r1 = cd.process_ref(ref, citation, doi_map.get(ref, {}), pdf_dir, 0.3)
                if r1.get("status") == "ok" or (os.path.exists(out_path) and os.path.getsize(out_path) > 5000):
                    if _verify_download(ref, citation, out_path):
                        print("    -> OK %s" % r1.get("source", "exists"), flush=True)
                        results.append({"ref": ref, "status": "ok", "source": r1.get("source")})
                        ok = True
                    else:
                        _safe_remove(out_path)
                        print("    -> 内容核验失败 (mismatch), 删除重试", flush=True)
            except Exception as e:
                print("    -> cascade 异常: %s" % str(e)[:80], flush=True)
        if not ok:
            try:
                r2 = rd.process_ref(ref, citation, doi_map.get(ref, {}), pdf_dir, 0.3)
                if r2.get("status") == "ok" or (os.path.exists(out_path) and os.path.getsize(out_path) > 5000):
                    if _verify_download(ref, citation, out_path):
                        print("    -> OK %s" % r2.get("source", "exists"), flush=True)
                        results.append({"ref": ref, "status": "ok", "source": r2.get("source")})
                    else:
                        _safe_remove(out_path)
                        print("    -> 内容核验失败 (mismatch), 链接保留表格", flush=True)
                        results.append({"ref": ref, "status": "failed"})
                else:
                    print("    -> FAILED (链接保留在表格 H 列)", flush=True)
            except Exception as e:
                print("    -> round2 异常: %s" % str(e)[:80], flush=True)
        if not ok:
            # PubMed 术语检索兜底 (正文句含英文医学术语时)
            try:
                if _pubmed_term_download(citation, out_path):
                    if _verify_download(ref, citation, out_path):
                        print("    -> OK pubmed-term", flush=True)
                        results.append({"ref": ref, "status": "ok", "source": "pubmed_term"})
                        ok = True
                    else:
                        _safe_remove(out_path)
            except Exception:
                pass
        if not ok:
            results.append({"ref": ref, "status": "failed"})
    elapsed = int(time.monotonic() - start)
    print("  下载阶段: %d s, 标号 ok=%d, 全文库 ok=%d" % (
        elapsed, sum(1 for r in results if r.get("status") == "ok"), full_ok), flush=True)
    report = {"budget_s": budget_s, "elapsed_s": elapsed, "budget_hit": budget_hit,
              "results": results, "full_lib_ok": full_ok}
    json.dump(report, open(os.path.join(root, "_download_auto_report.json"), "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    return report


# ============ Step 4: 整理下载目录 ============
def step_organize(root):
    """_2_pdfs 扁平 → _literature_citation_index/Pn-x/Pn-x_main.pdf"""
    src = os.path.join(root, "_2_pdfs")
    dst = os.path.join(root, "_literature_citation_index")
    n = 0
    if os.path.isdir(src):
        for f in sorted(os.listdir(src)):
            m = re.match(r"(P\d+-\d+)\.pdf$", f)
            if not m:
                continue
            pn = m.group(1)
            d = os.path.join(dst, pn)
            os.makedirs(d, exist_ok=True)
            dest = os.path.join(d, pn + "_main.pdf")
            if not (os.path.exists(dest) and os.path.getsize(dest) == os.path.getsize(os.path.join(src, f))):
                shutil.copy2(os.path.join(src, f), dest)
                n += 1
    print("[4] 整理下载目录: %d 个 Pn-x → %s" % (n, dst), flush=True)
    return n


# ============ Step 5: highlight plan ============
def step_plans(pptx, root):
    from tma_highlight_by_slide import extract_slide_visual, slide_terms
    from pptx import Presentation
    prs = Presentation(pptx)
    n_slides = len(prs.slides)
    vpath = os.path.join(root, "_vision_report.json")
    vision = json.load(open(vpath, encoding="utf-8")) if os.path.exists(vpath) else None
    pdf_dir = os.path.join(root, "_2_pdfs")
    pdfs = [f[:-4] for f in os.listdir(pdf_dir)] if os.path.isdir(pdf_dir) else []
    plans_dir = os.path.join(root, "_highlight_plans")
    os.makedirs(plans_dir, exist_ok=True)
    n = 0
    for slide_num in range(1, n_slides + 1):
        visual = extract_slide_visual(pptx, slide_num, vision)
        slide_pdfs = [p for p in pdfs if p.startswith("P%d-" % slide_num)]
        terms, data = slide_terms(visual)
        plan = {
            "slide": slide_num,
            "n_text_blocks": len(visual.get("text_blocks") or []),
            "n_tables": len(visual.get("tables") or []),
            "n_images": len(visual.get("images") or []),
            "citation_marks": [str(m.get("mark")) for m in (visual.get("citation_marks") or [])],
            "pdfs": sorted(slide_pdfs),
            "keywords": sorted(list(terms))[:60],
            "data_points": [d for d in data if d][:30],
            "status": "pending",
        }
        json.dump(plan, open(os.path.join(plans_dir, "slide_%03d_plan.json" % slide_num), "w", encoding="utf-8"),
                  ensure_ascii=False, indent=2)
        n += 1
    print("[5] highlight plan: %d 个 slide → %s" % (n, plans_dir), flush=True)
    return n


# ============ Step 6: highlight (子进程 by-slide 全量) ============
def step_highlight(root):
    script = os.path.join(HERE, "tma_highlight_by_slide.py")
    cmd = [sys.executable, script, "--project-dir", root, "--force"]
    print("[6] 按 slide 顺序 highlight (by-slide)...", flush=True)
    r = subprocess.run(cmd, env=dict(os.environ, TMA_PROJECT=root))
    print("  highlight 退出码: %d" % r.returncode, flush=True)
    return r.returncode == 0


# ============ Step 7: 报告 ============
def step_report(root):
    # 报告兼容: final_report 依赖 _references_FINAL.json; manual_list 依赖缺失清单 (缺则从 _refs.json 生成)
    fin = os.path.join(root, "_references_FINAL.json")
    if not os.path.exists(fin) and os.path.exists(os.path.join(root, "_refs.json")):
        shutil.copy2(os.path.join(root, "_refs.json"), fin)
    ml = os.path.join(root, "_manual_download_list.json")
    if not os.path.exists(ml):
        try:
            refs = json.load(open(os.path.join(root, "_refs.json"), encoding="utf-8"))
            pdf_dir = os.path.join(root, "_2_pdfs")
            have = set(f[:-4] for f in os.listdir(pdf_dir)) if os.path.isdir(pdf_dir) else set()
            miss = [{"ref": k, "citation": v} for k, v in sorted(refs.items())
                    if k not in have]
            json.dump(miss, open(ml, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
        except Exception as e:
            print("  [warn] 缺失清单生成失败: %s" % str(e)[:80], flush=True)
    for script in ["tma_final_report.py", "tma_manual_list.py"]:
        cmd = [sys.executable, os.path.join(HERE, script)]
        r = subprocess.run(cmd, env=dict(os.environ, TMA_PROJECT=root))
        print("  %s 退出码: %d" % (script, r.returncode), flush=True)
    # 全文库对照表: 参考文献列表 16 条完整引文 ↔ 下载状态 (供人工对照 PPT 标号)
    _write_full_lib_table(root)
    print("[7] 交付报告已生成 (CSV + 交付报告 + 人工下载清单 + 全文库对照表)", flush=True)


def _write_full_lib_table(root):
    try:
        full_path = os.path.join(root, "_references_full.json")
        if not os.path.exists(full_path):
            return
        full_refs = json.load(open(full_path, encoding="utf-8"))
        pdf_dir = os.path.join(root, "_2_pdfs")
        lines = ["# 全文库对照表 (参考文献列表 → 下载状态)", ""]
        lines.append("> PPT 正文上标标号为参考文献列表全局编号; 人工下载后放入 `_2_pdfs/ref{N}.pdf` 即可。")
        lines.append("")
        lines.append("| 编号 | 完整引文 | PDF |")
        lines.append("|---|---|---|")
        for num in sorted(full_refs, key=lambda x: int(x)):
            cit = str(full_refs[num])[:90]
            ref_pdf = os.path.join(pdf_dir, "ref%s.pdf" % num)
            ok = os.path.exists(ref_pdf) and os.path.getsize(ref_pdf) > 5000
            lines.append("| %s | %s | %s |" % (num, cit, "✅ ref%s.pdf" % num if ok else "❌ 待人工下载"))
        out = os.path.join(root, "_全文库对照表.md")
        with open(out, "w", encoding="utf-8") as f:
            f.write(chr(10).join(lines) + chr(10))
        print("  全文库对照表 → %s" % out, flush=True)
    except Exception as e:
        print("  [warn] 全文库对照表失败: %s" % str(e)[:80], flush=True)


# ============ 主流程 ============
def run_pipeline(nl_text=None, ppt=None, project_dir=None, budget_s=DEFAULT_BUDGET_S,
                 limit=0, skip_render=False, skip_download=False, skip_highlight=False,
                 skip_report=False):
    intent = parse_nl(nl_text) if nl_text else {
        "pptx": ppt, "render": True, "extract": True, "download": True, "highlight": True
    }
    pptx = intent.get("pptx") or ppt
    if not pptx or not os.path.isfile(pptx):
        print("错误: 找不到 PPT 文件: %s" % pptx)
        return 1
    root = resolve_project(pptx, project_dir)
    os.makedirs(root, exist_ok=True)
    print("== via54Medit 全自动管线 ==")
    print("  PPT: %s" % pptx)
    print("  项目根: %s" % root)
    print("  意图: 渲染=%s 下载=%s highlight=%s" % (intent["render"], intent["download"], intent["highlight"]))

    from deps_auto import ensure_env
    ok, problems = ensure_env(install=True)
    if not ok:
        print("⚠️ 环境自检有 %d 项问题, 继续尝试 (缺依赖的步骤会降级)" % len(problems), flush=True)

    # PPT 复制进项目根 (by-slide highlight 需项目内 PPTX 自动发现)
    proj_pptx = os.path.join(root, os.path.basename(pptx))
    if os.path.abspath(pptx) != os.path.abspath(proj_pptx):
        try:
            shutil.copy2(pptx, proj_pptx)
        except Exception as e:
            print("  [warn] PPT 复制进项目失败: %s" % str(e)[:80], flush=True)
    if intent["render"] and not skip_render:
        step_render(pptx, root)
    step_extract_refs(pptx, root)
    if intent["download"] and not skip_download:
        step_download(root, budget_s=budget_s, limit=limit)
    step_organize(root)
    if intent["highlight"] and not skip_highlight:
        step_plans(pptx, root)
        step_highlight(root)
    if not skip_report:
        step_report(root)
    print("\n✅ 全自动管线完成: %s" % root, flush=True)
    return 0


def main():
    parser = argparse.ArgumentParser(description="via54Medit 自然语言一键全自动管线")
    parser.add_argument("nl", nargs="?", help='自然语言指令, 如 "帮我识别 X.pptx 的文献引用，下载并highlight"')
    parser.add_argument("--ppt", default=None, help="PPT 路径 (或写在自然语言里)")
    parser.add_argument("--project-dir", default=None, help="项目根 (默认 PPT 所在目录)")
    parser.add_argument("--budget", type=int, default=DEFAULT_BUDGET_S, help="下载预算秒数 (默认 3600)")
    parser.add_argument("--limit", type=int, default=0, help="仅处理前 N 条待下载 (调试)")
    parser.add_argument("--skip-render", action="store_true")
    parser.add_argument("--skip-download", action="store_true")
    parser.add_argument("--skip-highlight", action="store_true")
    parser.add_argument("--skip-report", action="store_true")
    ns = parser.parse_args()
    if not ns.nl and not ns.ppt:
        parser.print_help()
        return 1
    return run_pipeline(nl_text=ns.nl, ppt=ns.ppt, project_dir=ns.project_dir,
                        budget_s=ns.budget, limit=ns.limit,
                        skip_render=ns.skip_render, skip_download=ns.skip_download,
                        skip_highlight=ns.skip_highlight, skip_report=ns.skip_report)


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    sys.exit(main())
