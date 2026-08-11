#!/usr/bin/env python3
"""
ppt_understand.py — Step 1: PPT 视觉理解 + 引文标号位置提取 (v2.0)

═══════════════════════════════════════════════════════════════════════════
硬规则 #1 (根因): 之前 find_citation_marks 用泛化正则把期刊年份/ISSN 等
数字都当成标号, 导致匹配到 20+ 个噪音标号. v2 改为语义驱动: 只从表格
"方案名称/药物"列和标题横幅中提取"中文/英文词+数字"的引用标号.

硬规则 #2 (GitHub 高 star): python-pptx (PyPI, 4k+/week) — PPTX 读写标准库
═══════════════════════════════════════════════════════════════════════════

核心能力:
  - extract_ppt_slide()          → PPT 第 N 页完整结构 (text, table, bbox)
  - find_citation_marks_v2()     → 语义驱动引文标号 (中文词+数字, 方案+数字)
  - get_ppt_mark_context()       → 单标号 PPT 上下文 (标号, 位置, 视觉, 文本)
  - build_ppt_vision_report()    → PPT 视觉理解报告 (Step 1 输出)

输出示例 (P5):
  标号 1: 标题横幅 "uHCC一线治疗方案1"
  标号 2: Row1 药物列 "奥沙利铂+5‑FU+亚叶酸钙2"
  标号 3: Row2 药物列 "索拉非尼3,4"
  标号 4: Row2 药物列 "索拉非尼3,4" (与 3 共享)
  ...
"""

import sys
import re
from pathlib import Path
from collections import OrderedDict
from typing import Dict, List, Optional, Any

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    Presentation = None

PPT_PATH = "/Users/david/Desktop/雷管方案_文献整理/PPT原版_雷管方案_三重获益_引领uHCC一线治疗_0622.pptx"


# ═══════════════════════════════════════════════════════════════════════════
# 核心函数
# ═══════════════════════════════════════════════════════════════════════════


def extract_ppt_slide(pptx_path: str, slide_num: int) -> List[Dict]:
    """
    提取 PPT 第 N 页的完整结构
    slide_num: 1-indexed
    返回每个 shape: {idx, name, type, text/coordinates/table_data}
    """
    if not HAS_PPTX:
        raise ImportError("python-pptx 未安装")

    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num - 1]
    structures = []

    for i, shape in enumerate(slide.shapes):
        item = {
            "idx": i,
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left": int(shape.left) if shape.left else 0,
            "top": int(shape.top) if shape.top else 0,
            "width": int(shape.width) if shape.width else 0,
            "height": int(shape.height) if shape.height else 0,
        }

        if shape.has_text_frame:
            item["type"] = "text"
            item["text_full"] = shape.text_frame.text
            paragraphs = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            item["text"] = "\n".join(paragraphs)

        elif shape.has_table:
            table = shape.table
            item["type"] = "table"
            rows_data = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows_data.append(cells)
            item["table_data"] = rows_data
            item["table_rows"] = len(rows_data)
            item["table_cols"] = len(rows_data[0]) if rows_data else 0
        else:
            item["type"] = str(shape.shape_type)

        structures.append(item)
    return structures


def _extract_table_citation_marks(table_data: List[List[str]]) -> Dict[int, Dict]:
    """
    从表格中提取引文标号 (v2 语义驱动)

    提取规则:
      1) 方案名称列 / 药物列中的"中文词+数字"模式
         - "仑伐替尼5" → 5
         - "T+A方案8,9" → 8, 9
         - "双达方案11" → 11
         - "雷管方案17,18" → 17, 18
      2) 药物列中"中文+数字"末尾
         - "奥沙利铂+5‑FU+亚叶酸钙2" → 2
      3) "O+Y15,16" 特殊模式 → 15, 16
    """
    marks = OrderedDict()

    for row_idx, row in enumerate(table_data):
        if row_idx == 0:  # header
            continue

        # 只从 col1(方案名称) 和 col2(药物) 提取
        col_map = [(1, "方案名称"), (2, "药物")]
        for col_idx, col_name in col_map:
            if col_idx >= len(row):
                continue
            cell = row[col_idx]
            if not cell or not cell.strip():
                continue

            found = set()

            # 规则 1a: 中文词 + 数字序列 (含逗号), 如 "索拉非尼3,4" → 3, 4
            for m in re.finditer(r'[\u4e00-\u9fff]+([0-9,]+)', cell):
                for n in re.findall(r'\d+', m.group(1)):
                    found.add(int(n))

            # 规则 1b: 中文词 + 数字 (不含逗号, 兜底)
            for m in re.finditer(r'[\u4e00-\u9fff]+([0-9]+)', cell):
                found.add(int(m.group(1)))

            # 规则 2: "方案" + 数字, 如 "T+A方案8,9" (中文词没匹配到, 因为 T 是英文)
            for m in re.finditer(r'[方案]+([0-9,]+)', cell):
                for n in re.findall(r'\d+', m.group(1)):
                    found.add(int(n))

            # 规则 3: "O+Y" + 数字, 如 "O+Y15,16"
            for m in re.finditer(r'[A-Z]+\+Y?\s*([0-9,]+)', cell):
                for n in re.findall(r'\d+', m.group(1)):
                    found.add(int(n))

            for num in sorted(found):
                if num not in marks:
                    marks[num] = {
                        "row": row_idx,
                        "column": col_name,
                        "context": cell.strip(),
                        "row_context": " | ".join(row),
                    }

    return marks


def find_citation_marks_v2(
    pptx_path: str, slide_num: int
) -> Dict[int, Dict]:
    """
    语义驱动引文标号提取 (v2 + v3 增强)

    步骤:
      1. 从所有表格中提取 (规则 1-3, 适用于 slide 5 表格)
      2. 从标题横幅中提取 (如 "uHCC一线治疗方案1" → 标号1)
      3. 从正文中提取 (适用于 slide 3/4 自由布局文本)
         - "中国肝癌新发和死亡病例占全球近半数1" → 1
         - "中国肝癌5年生存率仅14.4%, 远低于其他癌种3" → 3
         - "《健康中国行动...》2" → 2
         - "中晚期肝癌患者 显著拉低总体生存率4" → 4

    排除: 期刊年份 (JCO 2013 → 2013), ISSN, 疗效数据 (mOS 12.1月)
    """
    structures = extract_ppt_slide(pptx_path, slide_num)
    marks = OrderedDict()

    # 1. 表格提取
    for item in structures:
        if item["type"] == "table":
            table_marks = _extract_table_citation_marks(item["table_data"])
            for num, info in table_marks.items():
                marks[num] = {
                    **info,
                    "shape_idx": item["idx"],
                    "shape_name": item["name"],
                    "extraction": "table",
                }

    # 2. 标题横幅提取: "uHCC一线治疗方案1" → 1
    # 3. 正文文本提取: "中国肝癌新发和死亡病例占全球近半数1" → 1
    for item in structures:
        if item["type"] == "text":
            text = item.get("text_full", "")
            if not text:
                continue

            # 匹配: 每行末尾 "中文词/标点 + 数字"
            # 支持: 单标号 (其他癌种3 → 3)
            #       逗号分隔 (中国肝癌HBV感染率高1,2 → 1, 2)
            #       连字符范围 (中国肝癌患者肿瘤负荷大1-3 → 1, 2, 3)
            #       书名号/中文标点 (《健康中国行动...》2 → 2)
            for line in text.split(chr(10)):
                line = line.strip()
                if not line:
                    continue

                # v9.5: 过滤纯数字 (图表刻度/数据标签, e.g. "20", "40", "24")
                if re.fullmatch(r'\d+\.?\d*', line):
                    continue

                # v9.5: 过滤医学缩写+数字 结尾 (PD-L1, CTLA-4, mOS, ORR 等术语)
                # 兼容中文标点 (：:) 和描述文字 (细胞毒性T淋巴细胞相关蛋白)
                medical_term_suffix = re.search(
                    r'(?:PD-?L1|PD-?1|CTLA-?4|mOS|mPFS|mTTP|mTTR|ORR|DCR|PFS|TTP|TTR|HCC|uHCC|HBV|HCV|AFP|ALT|AST|TBIL|ALBI|ECOG|BOR|CR|PR|SD|PD|AE|TRAE|AESI|HR|CI|OS|STRIDE|T\+A|T\+D|O\+Y|LEN|SOR|NIVO|IPI|Cabo|Regor|Atezo|Bev|Camrelizumab|Atezolizumab|Durvalumab|Tremelimumab|Tislelizumab|Lenvatinib|Sorafenib|Cabozantinib|Apatinib|Sintilimab|NATREC|CAR|TACE|HIFU|RFA|MWA).*?\d+\s*$',
                    line, re.IGNORECASE
                )
                if medical_term_suffix:
                    # 进一步: 如果 line 是医学缩写解释段 (含 2+ 医学缩写 + 中文描述)
                    # 排除 PD-L1/CTLA-4 解释段 (PD-L1: ... PD-L1: ... CTLA-4: ...)
                    has_terminology_explanation = (
                        line.count('PD-L1') >= 2 or
                        (line.count('PD-L1') >= 1 and line.count('CTLA-4') >= 1) or
                        line.count('程序性死亡配体') >= 1 or
                        line.count('细胞毒性T淋巴细胞') >= 1
                    )
                    if has_terminology_explanation:
                        continue
                    # 其他: 只过滤明显的医学术语结尾 (PD-L1, CTLA-4 在末尾)
                    if re.search(
                        r'(?:PD-?L1|PD-?1|CTLA-?4|HCC|uHCC|HBV|HCV|AFP|mOS|mPFS|ORR|DCR|PFS|TTP|TTR|TRAE|STRIDE|NIVO|IPI|OS|HR)\s*$',
                        line, re.IGNORECASE
                    ):
                        continue

                # v9.5: 过滤研究名 + 数字 (ORIENT-3, CheckMate-9DW, NCTxxxxxx 等)
                study_name_suffix = re.search(
                    r'(?:ORIENT|CARES|CHECKMATE|IMBRAVE|HIMALAYA|IMFINZI|KEYNOTE|JUPITER|COSMIC|SHARP|RESORCE|REFLECT|CELESTIAL|BRYTON|TACTICS|STRIDE|APASL|ASCO|ESMO|NCCN|AHCC|BMS|MRK)\s*[-]?\s*\d+\s*$',
                    line, re.IGNORECASE
                )
                if study_name_suffix:
                    continue

                # v9.5: 过滤 N=, OR=, CR=, AND= 等统计表达式中的数字
                stat_expr_suffix = re.search(
                    r'(?:N=|N≥|N≤|OR=|AND=|CR=|p=|P=|HR=|95%\s*CI|n=|CR:|p:|HR:|95%:|CI:)\s*[\d\.]+\s*$',
                    line, re.IGNORECASE
                )
                if stat_expr_suffix:
                    continue

                # v9.5: 过滤期刊卷号 (J Hematol Oncol. 17(1): 13 → 1 和 13)
                volume_page_suffix = re.search(
                    r'\d+\(\d+\)\s*[:：]\s*\d+\s*$',
                    line
                )
                if volume_page_suffix:
                    continue

                # v9.5: 过滤作者+年份 (Finn et al. New Engl J Med. 2020 → 20)
                author_year_suffix = re.search(
                    r'(?:et al\.?|J\.|J\s+)?(?:N\s+Engl|N\s+Engl\s+J\s+Med|Lancet|JAMA|Nat|J\s+Hepatol|J\s+Clin|JAMA\s+Oncol|Lancet\s+Oncol|Nat\s+Med|Nat\s+Commun|Nat\s+Rev|Cancer\s+Cell|Cancer\s+Discov|Cancer\s+Res|Clin\s+Cancer\s+Res|Gastroenterology|Hepatology|J\s+Immunother|Front\s+Oncol|Med|Br\s+J\s+Cancer|Ann\s+Oncol|Eur\s+J\s+Cancer)\s*\.?\s*\d{4}\s*$',
                    line, re.IGNORECASE
                )
                if author_year_suffix:
                    continue

                # v9.5: 过滤文献作者信息 (Zhang, Huajun et al. ...2022 → 22)
                author_text_suffix = re.search(
                    r'(?:[A-Z][a-z]+,\s*[A-Z][a-z]+|et al\.?)[^.]*\b(20\d{2}|19\d{2})\s*$',
                    line
                )
                if author_text_suffix:
                    continue

                range_m = re.search(r'([1-9][0-9]?)\s*-\s*([1-9][0-9]?)\s*$', line)
                if range_m:
                    start = int(range_m.group(1))
                    end = int(range_m.group(2))
                    if 1 <= start <= end <= 24 and end - start <= 5:
                        for num in range(start, end + 1):
                            if num not in marks:
                                marks[num] = {
                                    "context": line,
                                    "row": None,
                                    "column": None,
                                    "shape_idx": item["idx"],
                                    "shape_name": item["name"],
                                    "extraction": "text_range",
                                }
                        continue

                # 2) 单/逗号分隔 (其他癌种3 → 3 / HBV感染率高1,2 → 1,2)
                # 末尾连续数字+逗号+数字...
                seq_match = re.search(r'([1-9][0-9]?(?:,\s*[1-9][0-9]?)*)\s*$', line)
                if seq_match:
                    nums_str = seq_match.group(1)
                    nums = [int(n.strip()) for n in nums_str.split(',')]
                    for num in nums:
                        if 1 <= num <= 24 and num not in marks:
                            marks[num] = {
                                "context": line,
                                "row": None,
                                "column": None,
                                "shape_idx": item["idx"],
                                "shape_name": item["name"],
                                "extraction": "text_v3",
                            }
                    if nums:
                        continue

                # 3) 兜底: 末尾 单标号 (中文词+数字)
                matches = list(re.finditer(r'[\u4e00-\u9fff《》（）]+([1-9][0-9]?)$', line))
                if matches:
                    num = int(matches[-1].group(1))
                    if 1 <= num <= 24 and num not in marks:
                        marks[num] = {
                            "context": line,
                            "row": None,
                            "column": None,
                            "shape_idx": item["idx"],
                            "shape_name": item["name"],
                            "extraction": "text_v3",
                        }

    return marks


def get_ppt_mark_context(
    pptx_path: str, slide_num: int, mark_num: int
) -> Dict:
    """
    获取单个标号的完整 PPT 上下文
    """
    marks = find_citation_marks_v2(pptx_path, slide_num)
    if mark_num not in marks:
        return {"mark_num": mark_num, "found": False}
    return {"mark_num": mark_num, "found": True, **marks[mark_num]}


def build_ppt_vision_report(
    pptx_path: str, slide_num: int, expected_marks: List[int]
) -> Dict:
    """
    构建 PPT 视觉理解报告 (Step 1 输出)
    """
    structures = extract_ppt_slide(pptx_path, slide_num)
    marks = find_citation_marks_v2(pptx_path, slide_num)

    report = {
        "slide_num": slide_num,
        "total_shapes": len(structures),
        "title_text": "",
        "tables": [],
        "citation_marks": {},
        "missing_marks": [],
        "found_marks": sorted(marks.keys()),
    }

    # 提取标题
    for s in structures:
        if s["type"] == "text" and "uHCC" in s.get("text_full", ""):
            report["title_text"] = s["text"].split('\n')[0][:60]

    # 提取表格元数据
    for s in structures:
        if s["type"] == "table":
            report["tables"].append({
                "shape_idx": s["idx"],
                "name": s["name"],
                "rows": s["table_rows"],
                "cols": s["table_cols"],
            })

    # 匹配预期标号
    for expected in expected_marks:
        if expected in marks:
            report["citation_marks"][expected] = marks[expected]
        else:
            report["missing_marks"].append(expected)

    return report


# ═══════════════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════════════


def cli_extract(args):
    slide_num = int(args[0]) if args else 5
    path = args[1] if len(args) > 1 else PPT_PATH
    structures = extract_ppt_slide(path, slide_num)
    print(f"✅ PPT 第 {slide_num} 页: {len(structures)} 个元素")
    for s in structures:
        if s["type"] == "table":
            print(f"  [{s['idx']}] {s['name']}: TABLE {s['table_rows']}x{s['table_cols']}")
        elif s["type"] == "text":
            txt = s.get("text", "")[:100]
            if txt.strip():
                print(f"  [{s['idx']}] {s['name']}: {txt}")
        else:
            print(f"  [{s['idx']}] {s['name']}: {s['type']}")


def cli_marks(args):
    slide_num = int(args[0]) if args else 5
    path = args[1] if len(args) > 1 else PPT_PATH
    marks = find_citation_marks_v2(path, slide_num)
    print(f"✅ v2 找到 {len(marks)} 个引文标号: {sorted(marks.keys())}\n")
    for num in sorted(marks.keys()):
        info = marks[num]
        ctx = info.get("context", "")[:80]
        row = info.get("row")
        col = info.get("column")
        shape = info.get("shape_name", "")
        loc = f"Row {row} / {col}" if row is not None else shape
        print(f"  标号 {num}: [{loc}] \"{ctx}\"")


def cli_context(args):
    mark_num = int(args[0]) if args else 1
    slide_num = int(args[1]) if len(args) > 1 else 5
    path = args[2] if len(args) > 2 else PPT_PATH
    ctx = get_ppt_mark_context(path, slide_num, mark_num)
    import json
    print(json.dumps(ctx, ensure_ascii=False, indent=2, default=str))


def main():
    if len(sys.argv) < 2:
        print(f"""用法: {sys.argv[0]} <command> [args]

命令:
  extract <slide_num> [pptx_path]      — 提取 PPT 第 N 页结构
  marks <slide_num> [pptx_path]        — v2 引文标号 (语义驱动)
  context <mark_num> [slide_num] [path] — 单标号上下文

示例:
  {sys.argv[0]} marks 5        — P5 引文标号 (16个)
  {sys.argv[0]} context 5 5    — P5 标号 5 上下文
  {sys.argv[0]} extract 5      — P5 结构
""")
        sys.exit(1)

    cmd = sys.argv[1]
    args = sys.argv[2:]
    cmds = {"extract": cli_extract, "marks": cli_marks, "context": cli_context}
    if cmd not in cmds:
        print(f"未知命令: {cmd}")
        sys.exit(1)
    cmds[cmd](args)


if __name__ == "__main__":
    main()